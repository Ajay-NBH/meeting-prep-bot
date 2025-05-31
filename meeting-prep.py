import datetime
import os.path
import time
import base64
from email.mime.text import MIMEText
import io # For GDrive downloads
import re
import markdown 
import json
import fitz

from google.auth.transport.requests import Request
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from googleapiclient.http import MediaIoBaseDownload


import google.generativeai as genai

# For parsing Office documents if downloaded from Drive
from pptx import Presentation
import openpyxl


# --- Configuration ---
# For Google Workspace APIs (Calendar, Gmail, Drive)
SCOPES = [
    'https://www.googleapis.com/auth/calendar.readonly', 
    'https://www.googleapis.com/auth/gmail.send',
    'https://www.googleapis.com/auth/drive.readonly',
    'https://www.googleapis.com/auth/spreadsheets.readonly', # For reading Google Sheets
    'https://www.googleapis.com/auth/calendar.events' # For tagging events - RECOMMENDED
]
CREDENTIALS_FILE = 'credentials.json' # Downloaded from GCP
TOKEN_FILE_PREFIX = 'token_brandvmeet' # Will generate token_brandvmeet_calendar.json etc.

# For Gemini API
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY") # Set this environment variable

# Google Drive Folder ID containing NBH data
NBH_GDRIVE_FOLDER_ID = os.getenv("NBH_GDRIVE_FOLDER_ID", "1rikXDq-ZyuZpUbN-ZLCsmcVJCswIlPDq") # Set env var or replace placeholder

AGENT_EMAIL = "brand.vmeet@nobroker.in" # Email of the agent account
ADMIN_EMAIL_FOR_NOTIFICATIONS = "ajay.saini@nobroker.in" # REPLACE with your actual email

PROCESSED_EVENTS_FILE = 'processed_event_ids.txt' # Simple file-based tracking for local runs

# Cache file for inferred industries
# INFERRED_INDUSTRIES_CACHE_FILE = 'inferred_industries_cache.json' # Cache for brand industry inference



# Specific File Names (you might make these configurable or discover them)
FILE_NAME_PITCH_DECK_PDF = "NBH Monetization Pitch Deck.pdf" # From your image
# FILE_NAME_CASE_STUDIES_GSLIDES = "National Campaigns_case studies" # From your image
FILE_NAME_CASE_STUDIES_PDF = "National_Campaigns_case_studies.pdf" # New or alternative
FILE_NAME_PHYSICAL_CAMPAIGNS_GSHEET = "Physical_campaigns_live_sheet" # From your image
FILE_NAME_DIGITAL_CAMPAIGNS_GSHEET = "Digital_Campaigns_live_sheet" # From your image (note "live sheet")
FILE_NAME_COM_DATA_GSHEET = "NoBroker_Overall_Data" # From your image
FILE_NAME_NBH_PREVIOUS_MEETINGS_GSHEET = "NBH_Previous_meetings"


def parse_names_from_cell_helper(cell_value_str):
    names = set()
    if not cell_value_str or str(cell_value_str).strip().lower() == 'n/a':
        return names
    
    cell_value_str = str(cell_value_str) # Ensure it's a string
    # Remove content in parentheses (e.g., "(NoBrokerHood)", "(Brand Representative)")
    cleaned_cell = re.sub(r'\s*\([^)]*\)', '', cell_value_str)
    cleaned_cell = cleaned_cell.replace('*', '').strip() # Remove asterisks and leading/trailing spaces

    # Split by common delimiters. \s* around delimiters handles spaces.
    potential_names = re.split(r'\s*[,;/&\n]\s*|\s+\band\b\s+|\s+\bwith\b\s+', cleaned_cell)
    
    for name_part in potential_names:
        final_name = name_part.strip().lower()
        # Filter for meaningful names (e.g., length > 2) and exclude common role descriptors
        if final_name and len(final_name) > 2 and \
           "nbh sales" not in final_name and \
           "brand representative" not in final_name and \
           "nobrokerhood" not in final_name and \
           "stay vista" not in final_name: # Example: add other generic terms to exclude
            names.add(final_name)
    return names

# --- Google Authentication and Service Building ---
def get_google_service(service_name, version, scopes_list, token_filename):
    creds = None
    if os.path.exists(token_filename):
        creds = Credentials.from_authorized_user_file(token_filename, scopes_list)
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            try:
                creds.refresh(Request())
            except Exception as e:
                print(f"Error refreshing token for {service_name}: {e}")
                creds = None # Force re-authentication
        if not creds: # Either no token or refresh failed
            flow = InstalledAppFlow.from_client_secrets_file(CREDENTIALS_FILE, scopes_list)
            creds = flow.run_local_server(port=0)
        with open(token_filename, 'w') as token:
            token.write(creds.to_json())
    try:
        service = build(service_name, version, credentials=creds)
        print(f"{service_name.capitalize()} service initialized successfully.")
        return service
    except HttpError as error:
        print(f'An error occurred building {service_name} service: {error}')
        return None
    
# --- Google Drive Functions ---
def list_files_in_gdrive_folder(drive_service, folder_id):
    # ... (Implementation from previous thought block - list files) ...
    # Ensure it handles empty folder_id
    if not folder_id or folder_id == "YOUR_GDRIVE_FOLDER_ID_HERE":
        print("Google Drive Folder ID for NBH data is not configured.")
        return []
    try:
        results = drive_service.files().list(
            q=f"'{folder_id}' in parents and trashed=false",
            # Include mimeType to decide how to read
            fields="nextPageToken, files(id, name, mimeType)"
        ).execute()
        items = results.get('files', [])
        return items
    except HttpError as error:
        print(f"An error occurred listing GDrive files for folder {folder_id}: {error}")
        return []

    

# THIS IS THE PRIMARY FUNCTION TO GET FILE DATA
def get_structured_gdrive_file_data(drive_service, sheets_service, file_id, file_name, mime_type):
    """
    Reads content from various file types in Google Drive.
    For PPTX (Google Slides export or native PPTX), returns a list of {"slide_number": X, "text": "content"} dicts.
    For Google Sheets or XLSX, returns a list of {"sheet_name": S, "row_index": R, "header": H, "values": V} dicts.
    For other text-based files, returns a single string of content.
    Returns a string error message if parsing fails for a specific type.
    """
    print(f"    Attempting to read structured data for: {file_name} (MIME: {mime_type})")

    MIMETYPE_GOOGLE_SHEET = 'application/vnd.google-apps.spreadsheet'
    MIMETYPE_GOOGLE_DOC = 'application/vnd.google-apps.document'
    MIMETYPE_GOOGLE_PRESENTATION = 'application/vnd.google-apps.presentation'
    MIMETYPE_MS_EXCEL = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    MIMETYPE_MS_POWERPOINT = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
    MIMETYPE_PDF = 'application/pdf'

    try:
        # --- Google Slides (Presentations) ---
        if mime_type == MIMETYPE_GOOGLE_PRESENTATION:
            print(f"    Exporting Google Slides '{file_name}' as PPTX for slide-level parsing...")
            # Export Google Slides as PPTX
            request_pptx = drive_service.files().export_media(
                fileId=file_id,
                mimeType=MIMETYPE_MS_POWERPOINT # Export as PPTX
            )
            fh_pptx = io.BytesIO()
            downloader_pptx = MediaIoBaseDownload(fh_pptx, request_pptx)
            done_pptx = False
            while not done_pptx:
                status, done_pptx = downloader_pptx.next_chunk()
            fh_pptx.seek(0)

            try:
                prs = Presentation(fh_pptx)
                slides_data = []
                for i, slide in enumerate(prs.slides):
                    slide_text_runs = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text_frame") and shape.text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if run.text:
                                        slide_text_runs.append(run.text)
                    slides_data.append({
                        "slide_number": i + 1,
                        "text": "\n".join(slide_text_runs),
                        "source_type": "Google Presentation (exported as PPTX)"
                    })
                if not slides_data and prs.slides: # Had slides but no text extracted
                     return f"Google Slides '{file_name}': Exported as PPTX, but no text content found in slides."
                elif not prs.slides: # Exported PPTX was empty
                    return f"Google Slides '{file_name}': Exported as PPTX, but it contained no slides."
                return slides_data                          
            except Exception as e_pptx_parse:
                print(f"    Warning: Failed to parse Google Slides '{file_name}' exported as PPTX: {e_pptx_parse}. Falling back to plain text export.")
                # Fallback: Export as plain text if PPTX export/parse fails
                request_text = drive_service.files().export_media(fileId=file_id, mimeType='text/plain')
                fh_text = io.BytesIO()
                downloader_text = MediaIoBaseDownload(fh_text, request_text)
                done_text = False
                while not done_text:
                    status, done_text = downloader_text.next_chunk()
                fh_text.seek(0)
                full_text = fh_text.read().decode('utf-8', errors='replace')
                if not full_text.strip():
                    return f"Google Slides '{file_name}': Fallback to plain text export, but no content found."
                return [{"slide_number": 1, "text": full_text, "source_type": "Google Presentation (all text fallback)"}]


        # --- Microsoft PowerPoint (.pptx) ---
        elif mime_type == MIMETYPE_MS_POWERPOINT: # Handles native PPTX files
            request = drive_service.files().get_media(fileId=file_id)
            fh = io.BytesIO()
            downloader = MediaIoBaseDownload(fh, request)
            done = False
            while not done: status, done = downloader.next_chunk()
            fh.seek(0)
            try:
                prs = Presentation(fh)
                slides_data = []
                for i, slide in enumerate(prs.slides):
                    slide_text_runs = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text_frame") and shape.text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if run.text: slide_text_runs.append(run.text)
                    slides_data.append({"slide_number": i + 1, "text": "\n".join(slide_text_runs), "source_type": "PPTX (native)"}) # Clarified source
                if not slides_data and prs.slides:
                    return f"PPTX File '{file_name}': Contained slides, but no text content found in them."
                elif not prs.slides:
                    return f"PPTX File '{file_name}': Contained no slides."
                return slides_data
            except Exception as e:
                return f"Could not parse native PPTX content for '{file_name}': {e}"
            

        # --- Google Sheets & Microsoft Excel ---
        elif mime_type == MIMETYPE_GOOGLE_SHEET or mime_type == MIMETYPE_MS_EXCEL:
            all_rows_data = []
            try:
                if mime_type == MIMETYPE_GOOGLE_SHEET:
                    spreadsheet = sheets_service.spreadsheets().get(spreadsheetId=file_id).execute()
                    for sheet_meta in spreadsheet.get('sheets', []):
                        sheet_title = sheet_meta['properties']['title']
                        # Read a significant number of rows, e.g., 1000, and all columns up to ZZ
                        range_str = f"'{sheet_title}'!A1:ZZ1000" 
                        result = sheets_service.spreadsheets().values().get(
                            spreadsheetId=file_id, range=range_str
                        ).execute()
                        rows = result.get('values', [])
                        if rows:
                            header = rows[0]
                            for row_idx, row_values in enumerate(rows): # (row_idx is 0-based here)
                                all_rows_data.append({"sheet_name": sheet_title, 
                                                      "row_index": row_idx + 1, # 1-based for display
                                                      "header": header, 
                                                      "values": row_values,
                                                      "source_type": "Google Sheet"})
                elif mime_type == MIMETYPE_MS_EXCEL:
                    request = drive_service.files().get_media(fileId=file_id)
                    fh = io.BytesIO()
                    downloader = MediaIoBaseDownload(fh, request)
                    done = False
                    while not done: status, done = downloader.next_chunk()
                    fh.seek(0)
                    workbook = openpyxl.load_workbook(fh)
                    for sheet_title in workbook.sheetnames:
                        sheet = workbook[sheet_title]
                        header = [cell.value for cell in sheet[1]] # Assuming header is 1st row
                        for row_idx, row_obj in enumerate(sheet.iter_rows(min_row=1, max_col=50, max_row=1000, values_only=True)): # max_col to limit width
                            # row_idx is 0-based here from iter_rows(min_row=1)
                            all_rows_data.append({"sheet_name": sheet_title,
                                                  "row_index": row_idx + 1, # 1-based for display
                                                  "header": header, 
                                                  "values": list(row_obj), # Ensure it's a list
                                                  "source_type": "Excel Sheet"})
                if not all_rows_data: # If loops completed but no data (e.g. all sheets were empty)
                    return f"Spreadsheet file '{file_name}' ({mime_type}) processed, but no data rows found in any sheet."
                return all_rows_data
            except Exception as e:
                return f"Could not parse Spreadsheet content for '{file_name}' ({mime_type}): {e}"
        
        # --- Google Documents ---
        elif mime_type == MIMETYPE_GOOGLE_DOC:
            request = drive_service.files().export_media(fileId=file_id, mimeType='text/plain')
            fh = io.BytesIO()
            downloader = MediaIoBaseDownload(fh, request)
            done = False
            while not done: status, done = downloader.next_chunk()
            fh.seek(0)
            return fh.read().decode('utf-8', errors='replace') # Returns string

        # --- PDF Files ---
        elif mime_type == MIMETYPE_PDF:
            try:
                request = drive_service.files().get_media(fileId=file_id)
                fh = io.BytesIO()
                downloader = MediaIoBaseDownload(fh, request)
                done = False
                while not done:
                    status, done = downloader.next_chunk()
                fh.seek(0)
                pdf_text = ""
                with fitz.open(stream=fh.read(), filetype="pdf") as doc:
                    for page in doc:
                        pdf_text += page.get_text() + "\n"
                if not pdf_text.strip():
                    return f"PDF File: '{file_name}'. No extractable text found."
                return pdf_text
            except Exception as e:
                return f"Could not parse PDF content: {e}" 

        # --- Plain Text Files ---
        elif mime_type.startswith('text/'):
            request = drive_service.files().get_media(fileId=file_id)
            fh = io.BytesIO()
            downloader = MediaIoBaseDownload(fh, request)
            done = False
            while not done: status, done = downloader.next_chunk()
            fh.seek(0)
            return fh.read().decode('utf-8', errors='replace') # Returns string
            
        else:
            return f"File type ({mime_type}) for '{file_name}' not configured for structured data extraction. Its name and existence might be relevant." # Returns string

    except HttpError as error:
        return f"An HTTP error occurred reading GDrive file {file_name} (ID: {file_id}): {error}" # Returns string
    except Exception as e:
        return f"A general error occurred reading GDrive file {file_name} (ID: {file_id}): {e}" # Returns string

def get_brand_industry(brand_name, gemini_llm_model):

    """
    Uses Gemini LLM to infer the industry for a given brand name.
    Returns a string such as "FMCG", "Real Estate", "E-commerce", etc.
    """
    if not gemini_llm_model:
        print("Gemini LLM model not available for industry inference.")
        return "Unknown"

    prompt = (
        f"Given the brand name '{brand_name}', infer the most likely industry or sector it operates in. "
        "Respond with only the industry name (e.g., 'FMCG', 'Real Estate', 'E-commerce', 'Automotive', 'Education', etc.). "
        "If you are unsure, respond with 'Unknown Industry'."
    )
    try:
        response = gemini_llm_model.generate_content(prompt)
        if response and response.candidates and response.candidates[0].content.parts:
            industry = response.candidates[0].content.parts[0].text.strip()
            # Clean up response
            if not industry or "unknown" in industry.lower():
                return "Unknown"
            # Optionally, take only the first word/phrase if LLM returns a sentence
            return industry.split('\n')[0].split('.')[0].strip()
        else:
            return "Unknown"
    except Exception as e:
        print(f"Error inferring industry for '{brand_name}': {e}")
        return "Unknown"

def summarize_file_content_with_gemini(gemini_llm_model, file_name, mime_type, file_content):
    """
    Uses Gemini LLM to summarize the content of a file for inclusion in the meeting brief.
    """
    if not gemini_llm_model:
        return "Error: Gemini model not available for summarization."

    prompt = (
        f"Summarize the following content from the file '{file_name}' (type: {mime_type}) in 50-70 concise bullet points, "
        "focusing on key facts, data, or insights that would be useful for a marketing/sales meeting. "
        "Do not include generic statements. If the content is not relevant, say 'No relevant content found.'\n\n"
        f"---\n{file_content}\n---"
    )
    try:
        response = gemini_llm_model.generate_content(prompt)
        if response and response.candidates and response.candidates[0].content.parts:
            summary = response.candidates[0].content.parts[0].text.strip()
            return summary
        else:
            return "No summary generated."
    except Exception as e:
        print(f"Error during Gemini summarization: {e}")
        return f"Error: Exception during Gemini summarization: {e}"

# --- Modify get_internal_nbh_data_for_brand ---
def get_internal_nbh_data_for_brand(drive_service, sheets_service, gemini_llm_model, 
                                    current_target_brand_name, current_meeting_data):
    # ... (initial checks for services and folder ID remain) ...
    print(f"Fetching and processing internal NBH data for target brand '{current_target_brand_name}'...")
    all_files_in_folder = list_files_in_gdrive_folder(drive_service, NBH_GDRIVE_FOLDER_ID) # Ensure this is called
    
    if not all_files_in_folder: # Add check for empty list
        return "No files found in the NBH GDrive folder.\n"

    target_brand_industry = "Unknown"
    if gemini_llm_model: # Only infer if LLM is available
        print(f"  Inferring industry for current target brand: '{current_target_brand_name}' using LLM...")
        target_brand_industry = get_brand_industry(current_target_brand_name, gemini_llm_model) # Singular call
        # Optional: You could implement a small, separate cache here for target_brand_industry
        # if you expect to run briefs for the same new target brand multiple times.
    else:
        print(f"  LLM not available. Industry for '{current_target_brand_name}' will be 'Unknown'.")
    
    print(f"  Target Brand: '{current_target_brand_name}', Inferred Industry: '{target_brand_industry}'")

    # Get clean names from current meeting data
    def get_clean_names_from_attendee_list(attendee_list_param):
        names = set()
        for att in attendee_list_param:
            name_to_process = att.get('displayName', att.get('email', '').split('@')[0])
            cleaned_name = name_to_process.lower().strip()
            for single_name in re.split(r'[,;/&]|\band\b|\bwith\b', cleaned_name):
                final_name = single_name.strip()
                if final_name and len(final_name) > 1:
                    names.add(final_name)
        return names

    # Use current_meeting_data here
    current_nbh_attendees_list = current_meeting_data.get('nbh_attendees', [])
    current_brand_attendees_list = current_meeting_data.get('brand_attendees_info', [])

    current_nbh_attendee_names = get_clean_names_from_attendee_list(current_nbh_attendees_list)
    current_brand_attendee_names = get_clean_names_from_attendee_list(current_brand_attendees_list)

    EXCLUDED_NBH_PSEUDO_NAMES_FOR_FOLLOWUP = {
        AGENT_EMAIL.lower().split('@')[0],
        "pia.brand"
    }

    current_nbh_attendee_names_for_followup_check = {
    name for name in current_nbh_attendee_names if name not in EXCLUDED_NBH_PSEUDO_NAMES_FOR_FOLLOWUP
    }
    #print(f"  DEBUG Current NBH Names for Follow-up Check: {current_nbh_attendee_names_for_followup_check}")
    #print(f"  DEBUG Current Brand Names for Follow-up Check: {current_brand_attendee_names}")

    final_context_parts = []

    # --- Process each specific, known file ---
    for item in all_files_in_folder: # Main loop starts here
        file_name = item.get('name', 'Unknown File')
        file_id = item['id']
        mime_type = item.get('mimeType', '')
        print(f"  Considering file: {file_name} ({mime_type})")

        # Call the consolidated and correctly named function
        file_data_object = get_structured_gdrive_file_data(
            drive_service, sheets_service, file_id, file_name, mime_type
        )

        # In get_internal_nbh_data_for_brand, after calling get_structured_gdrive_file_data:
        print(f"    DEBUG: file_data_object for {file_name} (Type: {type(file_data_object)}):")
        if isinstance(file_data_object, str):
            print(f"      Content/Error: {file_data_object[:500]}{'...' if len(file_data_object) > 500 else ''}")
        elif isinstance(file_data_object, list) and file_data_object:
            print(f"      Content: List of {len(file_data_object)} items. First item: {str(file_data_object[0])[:200]}...")
        elif isinstance(file_data_object, list) and not file_data_object:
            print(f"      Content: Empty list.")
        # Then proceed with your if/elif for file types

        # 1. NBH Monetization Pitch Deck.pdf (General Pitch Context)
        if FILE_NAME_PITCH_DECK_PDF.lower() in file_name.lower() and mime_type == 'application/pdf':
            if isinstance(file_data_object, str): # PDF content or error string
                is_error_or_no_content = (
                    (file_data_object.startswith("PDF File:") and "No extractable text found." in file_data_object) or
                    file_data_object.startswith("Could not parse PDF content:") or
                    file_data_object.startswith("An HTTP error occurred") or # General errors
                    file_data_object.startswith("A general error occurred") or
                    not file_data_object.strip() # Empty content
                )

                if not is_error_or_no_content and gemini_llm_model:
                     summary_of_pdf_content = summarize_file_content_with_gemini(gemini_llm_model, file_name, mime_type, file_data_object[:20000])
                     final_context_parts.append(f"## General NBH Pitch Overview (from '{file_name}'):\n{summary_of_pdf_content}\n")
                elif not is_error_or_no_content: # Has text, but no LLM
                     final_context_parts.append(f"## General NBH Pitch Overview (from '{file_name}'):\n{file_data_object[:1000]}...\n(Full content available, LLM summarization skipped)\n")
                else: # It's an error string or "no content" message
                     final_context_parts.append(f"## General NBH Pitch Context (from '{file_name}'):\n{file_data_object}\n")
            continue

        # 2. National Campaigns_case studies (Now as PDF)
        elif FILE_NAME_CASE_STUDIES_PDF.lower() in file_name.lower() and mime_type == 'application/pdf':
            print(f"    Processing Case Studies PDF: {file_name}")
            if isinstance(file_data_object, str) and file_data_object.strip(): # file_data_object is the extracted PDF text
                # Check if it's an error message from PDF parsing
                is_error_or_no_content = (
                    (file_data_object.startswith("PDF File:") and "No extractable text found." in file_data_object) or
                    file_data_object.startswith("Could not parse PDF content:") or
                    file_data_object.startswith("An HTTP error occurred") or
                    file_data_object.startswith("A general error occurred")
                )
                if is_error_or_no_content:
                    final_context_parts.append(f"## Case Studies (from PDF: '{file_name}'):\n{file_data_object}\n") # Report error                
                else:
                    # Successfully got text from the PDF.
                    # The PDF text will likely be one large block.
                    # Your LLM prompt needs to be good at segmenting this.
                    # We can't easily identify "slides" from a PDF's raw text dump.
                    # So, we provide the whole text and instruct the LLM.
                    # You might want to truncate if the PDF text is excessively long. MAX_PDF_TEXT_LEN = 50000      
                     
                    truncated_pdf_text = file_data_object
                    if len(truncated_pdf_text) > 50000: # Example truncation
                        truncated_pdf_text = file_data_object[:50000] + "\n... [Case Study PDF content truncated]"
                        print(f"    Truncated case study PDF text for '{file_name}' to 50000 chars.")                          

                    case_study_context = (
                        f"## Case Studies (from PDF document: '{file_name}'):\n"
                        f"The following is the extracted text content from the PDF. "
                        f"Please analyze this text to identify distinct case studies. For each, try to find "
                        f"the Brand, Objective, Activities, Key Results, and explain its relevance to '{current_target_brand_name}' "
                        f"(Industry: '{target_brand_industry}'). Pay attention to headings, lists, and common case study structures.\n\n"
                        f"---\nBEGIN PDF CASE STUDY TEXT\n---\n"
                        f"{truncated_pdf_text}\n"
                        f"---\nEND PDF CASE STUDY TEXT\n---\n"
                    )
                    final_context_parts.append(case_study_context)
            else: # Should not happen if get_structured_gdrive_file_data returns string or error
                final_context_parts.append(f"## Case Studies (from PDF: '{file_name}'):\nNo text data extracted or unexpected format.\n")
            continue # Move to next file

        # 3. & 4. Physical & Digital Campaigns (GSheets - Extract relevant rows)
        # THIS IS THE CORRECTED INDENTATION LEVEL FOR THE ELIF
        elif (FILE_NAME_PHYSICAL_CAMPAIGNS_GSHEET.lower() in file_name.lower() or \
              FILE_NAME_DIGITAL_CAMPAIGNS_GSHEET.lower() in file_name.lower()) and \
             mime_type == 'application/vnd.google-apps.spreadsheet':
            
            #rows_data_from_sheet = read_gdrive_file_content_with_slides(drive_service, sheets_service, file_id, file_name, mime_type)
            
            if isinstance(file_data_object, list) and file_data_object:
                # ... (the rest of your detailed logic for processing campaign GSheet rows) ...
                # ... (finding brand name column, iterating rows, checking direct match, checking industry match, etc.) ...
                relevant_rows_for_llm = [f"## Historical Campaign Data (from '{file_name}'):\n"]
                rows_added_count = 0
                MAX_CAMPAIGN_ROWS_TO_INCLUDE = 15 

                # Assuming first item in file_data_object (if it's a list of rows) contains the header
                first_row_dict = file_data_object[0] if isinstance(file_data_object[0], dict) else {}
                header_values = first_row_dict.get("header", [])
                
                brand_name_col_idx = -1
                industry_col_idx = -1 # NEW: For reading "Industry" column from sheet                
                

                if header_values:
                    try:
                        lower_header = [str(h).strip().lower() for h in header_values]
                        brand_name_col_idx = lower_header.index("brand name")
                        industry_col_idx = lower_header.index("industry") # Ensure this column exists in your sheets!
                    except ValueError:
                        missing_cols = []
                        if "brand name" not in lower_header: missing_cols.append("'Brand Name'")
                        if "industry" not in lower_header: missing_cols.append("'Industry'")
                        print(f"    Warning: Required column(s) {', '.join(missing_cols)} not found in header of {file_name}. Skipping detailed campaign processing for this file.")
                        final_context_parts.append(f"## Historical Campaign Data (from '{file_name}'):\nRequired columns missing in sheet header. Cannot process for relevant campaigns.\n")
                        continue # Skip to next file if essential columns are missing

                for row_info in file_data_object:
                    if not isinstance(row_info, dict) or 'values' not in row_info or not row_info['values']: continue
                    if row_info.get("values") == header_values : # Skip header row
                        continue

                    row_values = row_info['values']
                    is_relevant_row = False
                    relevance_reason = ""
                    row_brand_name = ""
                    row_campaign_industry_from_sheet = "" # NEW


                    # Attempt to extract Brand Name from the current row    
                    if brand_name_col_idx != -1 and len(row_values) > brand_name_col_idx and row_values[brand_name_col_idx] is not None:
                        row_brand_name = str(row_values[brand_name_col_idx]).strip()
                        if current_target_brand_name.lower() == row_brand_name.lower():
                            is_relevant_row = True
                            relevance_reason = f"Direct match for '{current_target_brand_name}'."
                    
                    # Attempt to extract Industry from the current row
                    if industry_col_idx != -1 and len(row_values) > industry_col_idx and row_values[industry_col_idx] is not None:
                        row_campaign_industry_from_sheet = str(row_values[industry_col_idx]).strip()

                    # Relevance check based on industry
                    if not is_relevant_row and target_brand_industry != "Unknown" and row_campaign_industry_from_sheet: # Check if sheet industry is not empty    
                       if row_campaign_industry_from_sheet.lower() == target_brand_industry.lower():
                               is_relevant_row = True
                               relevance_reason = f"Industry match for '{target_brand_industry}' in row industry '{row_campaign_industry_from_sheet}'."


                    # ROW DEBUG print statements would go here, after all checks
                    #print(f"    ROW DEBUG: TargetBrand='{current_target_brand_name.lower()}', RowBrand='{row_brand_name.lower()}'")
                    #print(f"    ROW DEBUG: TargetIndustry='{target_brand_industry.lower()}', RowSheetIndustry='{row_campaign_industry_from_sheet.lower()}'")
                    #print(f"    ROW DEBUG: is_relevant_row after all checks: {is_relevant_row}, Reason: '{relevance_reason}'")


                    if is_relevant_row and rows_added_count < MAX_CAMPAIGN_ROWS_TO_INCLUDE:
                        print(f"      MATCH FOUND for row {row_info.get('row_index')}: {relevance_reason}")
                        # ... your existing logic to add to relevant_rows_for_llm ...
                        row_details = []
                        for h_idx, h_val in enumerate(header_values): # Use header_values from the top of this file's processing block
                            cell_val = str(row_values[h_idx]) if len(row_values) > h_idx else "N/A"
                            if cell_val.strip() and cell_val != "N/A": # Only add if there's a value
                                row_details.append(f"{h_val}: {cell_val.strip()}")
                        if row_details: # Only add if we have some details  
                            relevant_rows_for_llm.append(f"### Campaign Row {row_info.get('row_index', 'N/A')} (Relevance: {relevance_reason}):\n" + ", ".join(row_details) + "\n---\n")
                        rows_added_count +=1
                
                if rows_added_count > 0:
                    final_context_parts.append("".join(relevant_rows_for_llm))
                else:
                    final_context_parts.append(f"## Historical Campaign Data (from '{file_name}'):\nNo campaigns found directly matching '{current_target_brand_name}' or its inferred industry '{target_brand_industry}'.\n")

            elif isinstance(file_data_object, str): 
                final_context_parts.append(f"## Note on '{file_name}':\nError processing campaign data: {file_data_object}\n")
            else: 
                 final_context_parts.append(f"## Historical Campaign Data (from '{file_name}'):\nNo data rows found or file not in expected sheet format.\n")
            continue

        # 5. com data _ Cross vertical Services_ West (GSheet - Database Numbers/Leads - MULTI-TAB)
        elif FILE_NAME_COM_DATA_GSHEET.lower() in file_name.lower() and mime_type == 'application/vnd.google-apps.spreadsheet':
            com_data_structured_output = [
                f"## NoBroker.com Platform Metrics (from GSheet: '{file_name}'):\n"
                f"This Google Sheet contains multiple tabs, each with specific NoBroker platform metrics (e.g., lead volumes, user data, city-wise breakdowns, service-specific numbers). The agent preparing this brief should analyze the data from the relevant tabs described below to identify compelling metrics for '{current_target_brand_name}' (Industry: '{target_brand_industry}').\n\n"
            ]
            
            # `read_gdrive_file_content_with_slides` already returns a list of dicts, 
            # where each dict can represent a row from a specific sheet (tab) within the GSheet.
            # Each dict has: {"sheet_name": S, "row_index": R, "header": H, "values": V}
            if isinstance(file_data_object, list) and file_data_object and "sheet_name" in file_data_object[0]:

            #all_rows_from_com_data_sheet = read_gdrive_file_content_with_slides(
                drive_service, sheets_service, file_id, file_name, mime_type
            #)

            if isinstance(file_data_object, list) and file_data_object:
                processed_tabs = {} # Key: tab_name, Value: list of row strings

                for row_info in file_data_object:
                    if 'sheet_name' not in row_info or 'values' not in row_info or not row_info['values']:
                        continue # Skip malformed row_info

                    tab_name = row_info['sheet_name']
                    
                    if tab_name not in processed_tabs:
                        processed_tabs[tab_name] = []
                        # Add header to the tab's data if available and it's the first time we see this tab
                        if "header" in row_info and row_info["header"]:
                             processed_tabs[tab_name].append(f"  Columns: {', '.join(map(str, row_info['header']))}\n")
                    
                    # Limit the number of sample rows per tab
                    MAX_SAMPLE_ROWS_PER_TAB = 7 # Show up to 7 data rows per tab (plus header)
                    if len(processed_tabs[tab_name]) < (MAX_SAMPLE_ROWS_PER_TAB + 1): # +1 for potential header
                        # Only add the row if it's not identical to the header (simple check)
                        if not (row_info.get("row_index", 0) == 1 and row_info.get("values") == row_info.get("header")):
                            processed_tabs[tab_name].append(f"  Data: {', '.join(map(str, row_info['values']))}\n")
                
                # Now assemble the output for the LLM
                for tab_name, rows_as_strings in processed_tabs.items():
                    if rows_as_strings: # Only if we have content for this tab
                        com_data_structured_output.append(f"### Data from Tab: '{tab_name}'\n")
                        com_data_structured_output.extend(rows_as_strings)
                        if len(rows_as_strings) > MAX_SAMPLE_ROWS_PER_TAB : # Check if we hit the row limit
                             com_data_structured_output.append("  ... (more rows available in this tab)\n")
                        com_data_structured_output.append("\n") # Separator between tabs

            elif isinstance(file_data_object, str): # Error reading sheet
                com_data_structured_output.append(f"Error processing the '{file_name}' GSheet: {file_data_object}\n") # CORRECTED: use file_data_object
            else: # No data or unexpected format
                com_data_structured_output.append(f"No data or tabs could be extracted from '{file_name}', or it was not in the expected sheet format.\n")
            
            final_context_parts.append("".join(com_data_structured_output))
            continue 
        
        # 6. NBH Previous Meetings Data
        elif FILE_NAME_NBH_PREVIOUS_MEETINGS_GSHEET.lower() in file_name.lower() and \
             mime_type == 'application/vnd.google-apps.spreadsheet':
            print(f"    Processing NBH Previous Meetings Sheet: {file_name}")
            previous_meeting_notes_list = [] # Use a list to build strings

            if isinstance(file_data_object, list) and file_data_object:
                # --- Define expected header names for this sheet (adjust as per your sheet) ---
                # It's crucial these match your sheet's column headers exactly (case-insensitive)
                # Based on your sample output: 'Meeitng ID', 'Meeting Date', 'Brand Name', 'Meeting Title',
                # 'Industry Type', 'Brand Size', 'Meeting Type', 'Meeting Agenda', 'NBH Participants', 'Client Participants',
                # 'Key Discussion Points', 'Action Items/Next Steps', 'Is Follow Up?'
                # Let's assume a few key ones for matching:
                HEADER_PREV_BRAND_NAME = "brand name"
                HEADER_PREV_MEETING_DATE = "meeting date" # For sorting
                HEADER_PREV_CLIENT_PARTICIPANTS_NAMES = "client participants" # Column with client names
                HEADER_PREV_NBH_PARTICIPANTS_NAMES = "nbh participants"       # Column with NBH names
                HEADER_PREV_KEY_DISCUSSION = "key discussion points"
                HEADER_PREV_ACTION_ITEMS = "action items"
                #HEADER_PREV_IS_FOLLOW_UP_TAG = "is follow up?" # A column you might add to explicitly tag follow-ups
                HEADER_PREV_KEY_QUESTIONS = "key questions" # Adjust to your exact column name
                HEADER_PREV_BRAND_TRAITS = "brand traits"   # Adjust
                HEADER_PREV_CUSTOMER_NEEDS = "customer needs" # Adjust
                HEADER_PREV_CLIENT_PAIN_POINTS = "client pain points" # Adjust

                first_row_dict = file_data_object[0] if isinstance(file_data_object[0], dict) else {}
                header_values = first_row_dict.get("header", [])
                lower_header = [str(h).strip().lower() for h in header_values]

                prev_brand_col_idx = -1
                prev_date_col_idx = -1
                # Add more column indices as needed based on what you want to extract
                print(f"    DEBUG NBH_Previous_meetings - Parsed Headers (lower_header): {lower_header}") # <--- ADD THIS

                try:
                            print(f"Attempting to find: '{HEADER_PREV_BRAND_NAME}'")
                            prev_brand_col_idx = lower_header.index(HEADER_PREV_BRAND_NAME)
                            print(f"Found '{HEADER_PREV_BRAND_NAME}' at index {prev_brand_col_idx}")

                            print(f"Attempting to find: '{HEADER_PREV_MEETING_DATE}'")
                            prev_date_col_idx = lower_header.index(HEADER_PREV_MEETING_DATE)
                            print(f"Found '{HEADER_PREV_MEETING_DATE}' at index {prev_date_col_idx}")

                            print(f"Attempting to find: '{HEADER_PREV_CLIENT_PARTICIPANTS_NAMES}'")
                            client_participants_names_col_idx = lower_header.index(HEADER_PREV_CLIENT_PARTICIPANTS_NAMES)
                            print(f"Found '{HEADER_PREV_CLIENT_PARTICIPANTS_NAMES}' at index {client_participants_names_col_idx}")

                            print(f"Attempting to find: '{HEADER_PREV_KEY_DISCUSSION}'")
                            key_discussion_col_idx = lower_header.index(HEADER_PREV_KEY_DISCUSSION)
                            print(f"Found '{HEADER_PREV_KEY_DISCUSSION}' at index {key_discussion_col_idx}")

                            print(f"Attempting to find: '{HEADER_PREV_ACTION_ITEMS}'")
                            action_items_col_idx = lower_header.index(HEADER_PREV_ACTION_ITEMS)
                            print(f"Found '{HEADER_PREV_ACTION_ITEMS}' at index {action_items_col_idx}")

                            print(f"Attempting to find: '{HEADER_PREV_KEY_QUESTIONS}'")
                            key_questions_col_idx = lower_header.index(HEADER_PREV_KEY_QUESTIONS)
                            print(f"Found '{HEADER_PREV_KEY_QUESTIONS}' at index {key_questions_col_idx}")

                            print(f"Attempting to find: '{HEADER_PREV_BRAND_TRAITS}'")
                            brand_traits_col_idx = lower_header.index(HEADER_PREV_BRAND_TRAITS)
                            print(f"Found '{HEADER_PREV_BRAND_TRAITS}' at index {brand_traits_col_idx}")

                            print(f"Attempting to find: '{HEADER_PREV_CUSTOMER_NEEDS}'")
                            customer_needs_col_idx = lower_header.index(HEADER_PREV_CUSTOMER_NEEDS)
                            print(f"Found '{HEADER_PREV_CUSTOMER_NEEDS}' at index {customer_needs_col_idx}")

                            print(f"Attempting to find: '{HEADER_PREV_CLIENT_PAIN_POINTS}'")
                            client_pain_points_col_idx = lower_header.index(HEADER_PREV_CLIENT_PAIN_POINTS)
                            print(f"Found '{HEADER_PREV_CLIENT_PAIN_POINTS}' at index {client_pain_points_col_idx}")

                            print(f"Attempting to find: '{HEADER_PREV_NBH_PARTICIPANTS_NAMES}'")
                            nbh_participants_names_col_idx = lower_header.index(HEADER_PREV_NBH_PARTICIPANTS_NAMES)
                            print(f"Found '{HEADER_PREV_NBH_PARTICIPANTS_NAMES}' at index {nbh_participants_names_col_idx}")

                except ValueError as e_val: # Catch the ValueError and print its specific message
                    print(f"    VALUE ERROR while finding column index: {e_val}") # This will tell you *which string* was not found
                    print(f"    Warning: Essential columns for previous meeting analysis not found in '{file_name}'. Skipping.")
                    # ... your existing error handling in except block ...
                    previous_meeting_notes_list.append(f"Could not process previous meetings sheet due to missing essential columns (e.g., Brand Name, Meeting Date, Key Discussion Points).\n")
                    final_context_parts.append("".join(previous_meeting_notes_list))
                    # Re-initialize all indices to -1 within the except block if one fails, so subsequent code doesn't use stale (but valid looking) indices from a previous file
                    prev_brand_col_idx = -1
                    prev_date_col_idx = -1
                    client_participants_names_col_idx = -1
                    key_discussion_col_idx = -1
                    action_items_col_idx = -1
                    key_questions_col_idx = -1
                    brand_traits_col_idx = -1
                    customer_needs_col_idx = -1
                    client_pain_points_col_idx = -1
                    nbh_participants_names_col_idx = -1
                    continue # Skip further processing for this file

                matching_previous_meetings_details = []

                for row_info in file_data_object:
                    if not isinstance(row_info, dict) or 'values' not in row_info or not row_info['values']:
                        continue
                    if row_info.get("values") == header_values:  # Skip header row
                        continue
                    
                    row_values = row_info['values']
                    prev_meeting_brand_name_from_sheet = "" # Initialize for current row
                    
                    # Extract brand name from the previous meeting row
                    prev_meeting_brand_name = ""
                    if len(row_values) > prev_brand_col_idx and row_values[prev_brand_col_idx] is not None:
                        prev_meeting_brand_name_from_sheet = str(row_values[prev_brand_col_idx]).strip()

                    # --- START OF TARGETED DEBUG BLOCK ---
                    # Only print if the sheet brand name (lowercase) potentially contains the target brand name (lowercase)
                    # or vice-versa. This helps narrow down the debug output significantly.
                    # This condition is intentionally broad to catch near misses or partial matches for debugging.
                    target_brand_lower = current_target_brand_name.lower()
                    sheet_brand_lower = prev_meeting_brand_name_from_sheet.lower()

                    if target_brand_lower in sheet_brand_lower or \
                       (sheet_brand_lower and target_brand_lower.startswith(sheet_brand_lower)) or \
                       (target_brand_lower and sheet_brand_lower.startswith(target_brand_lower)) or \
                       (target_brand_lower == "chitale" and "chitale" in sheet_brand_lower) or \
                       (target_brand_lower == "giva" and "giva" in sheet_brand_lower) : # Add more specific checks if needed

                        print(f"    TARGETED_PREV_MTG_DEBUG (Row {row_info.get('row_index', 'N/A')} in sheet '{file_name}'):")
                        print(f"        Current Target Brand (Original) : '{current_target_brand_name}'")
                        print(f"        Sheet Brand Name (Original)   : '{prev_meeting_brand_name_from_sheet}'")
                        print(f"        --- For Comparison ---")
                        print(f"        Current Target Brand (lower)  : '{target_brand_lower}' (len {len(target_brand_lower)})")
                        print(f"        Sheet Brand Name (lower)    : '{sheet_brand_lower}' (len {len(sheet_brand_lower)})")
                        # Optional: Byte representation for very tricky cases
                        # print(f"        Bytes Current Target (lower): {target_brand_lower.encode('utf-8', 'surrogateescape')}")
                        # print(f"        Bytes Sheet Brand (lower)   : {sheet_brand_lower.encode('utf-8', 'surrogateescape')}")
                        is_match = (sheet_brand_lower == target_brand_lower)
                        print(f"        Exact lowercase match?        : {is_match}")
                        print(f"        ------------------------------------")
                    # --- END OF TARGETED DEBUG BLOCK ---

                    if prev_meeting_brand_name_from_sheet.lower() == current_target_brand_name.lower():
                        # This previous meeting was with the same brand
                        print(f"      MATCH FOUND for previous meeting row {row_info.get('row_index')}: Brand match for '{current_target_brand_name}'")
                        meeting_details_dict = {
                            "original_row_index": row_info.get("row_index", "N/A") # Get it here
                        }    
                        meeting_details_dict = {"original_row_info": row_info} # Store original for later full extraction
                        try:
                            # Extract date for sorting (handle potential date parsing errors)
                            date_str = str(row_values[prev_date_col_idx]) if len(row_values) > prev_date_col_idx else None
                            if date_str:
                                # Attempt to parse common date formats, be robust
                                for fmt in ("%m/%d/%Y", "%Y-%m-%d", "%d/%m/%Y", "%d-%m-%Y", "%m/%d/%y"): # Add more if needed
                                    try:
                                        meeting_details_dict["date_obj"] = datetime.datetime.strptime(date_str.split(" ")[0], fmt).date() # Take only date part
                                        break
                                    except ValueError:
                                        pass
                                if "date_obj" not in meeting_details_dict:
                                     meeting_details_dict["date_obj"] = datetime.date.min # Fallback for unparseable dates
                            else:
                                meeting_details_dict["date_obj"] = datetime.date.min
                        except Exception:
                            meeting_details_dict["date_obj"] = datetime.date.min # Fallback

                        # Textual data
                        def get_cell_value(col_idx, default="N/A"):
                            if col_idx != -1 and len(row_values) > col_idx and row_values[col_idx] is not None and str(row_values[col_idx]).strip():
                                return str(row_values[col_idx]).strip()
                            return default

                        # Basic extraction for now, LLM can do more
                        meeting_details_dict["discussion"] = str(row_values[key_discussion_col_idx]) if len(row_values) > key_discussion_col_idx else "N/A"
                        meeting_details_dict["actions"] = str(row_values[action_items_col_idx]) if len(row_values) > action_items_col_idx else "N/A"
                        # New data points
                        meeting_details_dict["key_questions"] = str(row_values[key_questions_col_idx]) if len(row_values) > key_questions_col_idx and row_values[key_questions_col_idx] else "N/A"
                        meeting_details_dict["brand_traits"] = str(row_values[brand_traits_col_idx]) if len(row_values) > brand_traits_col_idx and row_values[brand_traits_col_idx] else "N/A"
                        meeting_details_dict["customer_needs"] = str(row_values[customer_needs_col_idx]) if len(row_values) > customer_needs_col_idx and row_values[customer_needs_col_idx] else "N/A"
                        meeting_details_dict["client_pain_points"] = str(row_values[client_pain_points_col_idx]) if len(row_values) > client_pain_points_col_idx and row_values[client_pain_points_col_idx] else "N/A"                

                        # Follow-up Check by Name
                        meeting_details_dict["is_direct_follow_up_candidate"] = False
                        prev_client_names_str = get_cell_value(client_participants_names_col_idx, "")
                        prev_nbh_names_str = get_cell_value(nbh_participants_names_col_idx, "")

                        prev_client_attendee_names = parse_names_from_cell_helper(prev_client_names_str)
                        prev_nbh_attendee_names_from_sheet = parse_names_from_cell_helper(prev_nbh_names_str)
                        prev_nbh_attendee_names_for_followup_check = {name for name in prev_nbh_attendee_names_from_sheet if name not in EXCLUDED_NBH_PSEUDO_NAMES_FOR_FOLLOWUP}
                        
                        common_brand_attendees_names = current_brand_attendee_names.intersection(prev_client_attendee_names)
                        common_nbh_attendees_names = current_nbh_attendee_names_for_followup_check.intersection(prev_nbh_attendee_names_for_followup_check)

                        if common_brand_attendees_names and common_nbh_attendees_names:
                            meeting_details_dict["is_direct_follow_up_candidate"] = True

                        matching_previous_meetings_details.append(meeting_details_dict)

                # --- Format the collected previous meeting data for the LLM ---
                previous_meeting_notes_list.append(f"## Insights from Previous NBH Meetings with '{current_target_brand_name}' (from '{file_name}'):\n")
                if matching_previous_meetings_details:
                    matching_previous_meetings_details.sort(key=lambda x: x["date_obj"], reverse=True)
                    MAX_PREVIOUS_MEETINGS_TO_NOTE = 3
                    
                    is_overall_direct_follow_up = any(mtg.get("is_direct_follow_up_candidate", False) for mtg in matching_previous_meetings_details[:MAX_PREVIOUS_MEETINGS_TO_NOTE])

                    if is_overall_direct_follow_up:
                        previous_meeting_notes_list.append(
                            f"**This upcoming meeting appears to be a DIRECT FOLLOW-UP based on attendee overlap with past meeting(s).** "
                            f"The brief should heavily focus on continuity, previous discussions, and action items.\n\n"
                        )
                    else:
                        previous_meeting_notes_list.append(
                            f"NBH has had previous interactions with '{current_target_brand_name}'. Key points are summarized below.\n\n"
                        )

                    for mtg_data in matching_previous_meetings_details[:MAX_PREVIOUS_MEETINGS_TO_NOTE]:
                        date_display = mtg_data["date_obj"].strftime("%Y-%m-%d") if mtg_data["date_obj"] != datetime.date.min else "Date N/A"
                        
                        # Access the 'row_index' from the nested 'original_row_info' dictionary
                        original_row_index_val = mtg_data.get("original_row_info", {}).get("row_index", "N/A")
                        note_parts = [f"### Previous Meeting on {date_display} (Row: {original_row_index_val})"]

                        if mtg_data.get("is_direct_follow_up_candidate"): note_parts.append(" *(Potential Direct Follow-up)*\n")
                        else: note_parts.append("\n")

                        if mtg_data['discussion'] != "N/A": note_parts.append(f"- **Key Discussion Points:** {mtg_data['discussion']}\n")
                        if mtg_data['key_questions'] != "N/A": note_parts.append(f"- **Key Questions Raised by Client:** {mtg_data['key_questions']}\n")
                        if mtg_data['brand_traits'] != "N/A": note_parts.append(f"- **Observed Brand Traits:** {mtg_data['brand_traits']}\n")
                        if mtg_data['customer_needs'] != "N/A": note_parts.append(f"- **Identified Customer Needs:** {mtg_data['customer_needs']}\n")
                        if mtg_data['client_pain_points'] != "N/A": note_parts.append(f"- **Client Pain Points Discussed:** {mtg_data['client_pain_points']}\n")
                        
                        # Conditionally include action items based on follow-up status
                        if mtg_data.get("is_direct_follow_up_candidate", False) and mtg_data['actions'] != "N/A":
                            note_parts.append(f"- **Action Items (Relevant for Follow-up):** {mtg_data['actions']}\n")
                        elif mtg_data['actions'] != "N/A": # If not a direct follow-up, but actions exist
                             note_parts.append(f"- **Action Items (General Context):** {mtg_data['actions']}\n")
                        
                        note_parts.append("---\n")
                        previous_meeting_notes_list.append("".join(note_parts))

                else:
                    previous_meeting_notes_list.append(f"No previous meeting records found specifically for '{current_target_brand_name}'.\n")
                
                final_context_parts.append("".join(previous_meeting_notes_list))

            elif isinstance(file_data_object, str): # Error reading sheet
                final_context_parts.append(f"## Previous NBH Meetings (from '{file_name}'):\nError processing the sheet: {file_data_object}\n")
            else: # No data
                final_context_parts.append(f"## Previous NBH Meetings (from '{file_name}'):\nNo data rows found in the sheet.\n")
            continue        

        # Fallback for other files or unrecognized structures
        elif isinstance(file_data_object, str) and file_data_object.strip():
            # Could be a simple text file, or an error message we haven't specifically handled
            # You might want to summarize these too if they are documents
            if "File type" in file_data_object or "Error" in file_data_object or "Could not parse" in file_data_object:
                 final_context_parts.append(f"## Note on file '{file_name}':\n{file_data_object}\n")
            elif gemini_llm_model: # It's some text content
                 summary = summarize_file_content_with_gemini(gemini_llm_model, file_name, mime_type, file_data_object[:15000])
                 final_context_parts.append(f"## Information from '{file_name}':\n{summary}\n")
            else: # No LLM, just pass truncated raw text
                 final_context_parts.append(f"## Information from '{file_name}':\n{file_data_object[:500]}...\n")# --- Other files (not specifically handled) ---

    # --- Combine all collected parts ---
    # ... (rest of the function to combine final_context_parts and truncate if necessary) ...
    if not final_context_parts:
        return "No specifically relevant internal NBH data could be extracted for this brand.\n"
    
    final_output_str = "\n\n".join(final_context_parts)
    
    MAX_TOTAL_WORDS_FOR_LLM_CONTEXT = 45000 
    current_word_count = len(final_output_str.split())
    if current_word_count > MAX_TOTAL_WORDS_FOR_LLM_CONTEXT:
        print(f"WARN: Combined internal data section is very long ({current_word_count} words). Truncating overall.")
        char_limit = MAX_TOTAL_WORDS_FOR_LLM_CONTEXT * 5 
        if len(final_output_str) > char_limit:
            cut_off_point = final_output_str.rfind(' ', 0, char_limit)
            if cut_off_point == -1: cut_off_point = char_limit
            final_output_str = final_output_str[:cut_off_point] + "\n... [Overall internal data truncated due to length]"

    return f"--- Targeted Internal NBH Data & Summaries for {current_target_brand_name} ---\n\n{final_output_str}\n--- End of Targeted Internal NBH Data & Summaries ---\n"


# --- Calendar Processing ---
def get_upcoming_meetings(calendar_service, calendar_id='primary', time_delta_hours=96): # Process meetings in next 3 days
    now_utc = datetime.datetime.utcnow()
    time_min_str = now_utc.isoformat() + 'Z'
    time_max_str = (now_utc + datetime.timedelta(hours=time_delta_hours)).isoformat() + 'Z'
    
    print(f'Getting events between {time_min_str} and {time_max_str}')
    try:
        events_result = calendar_service.events().list(
            calendarId=calendar_id, timeMin=time_min_str, timeMax=time_max_str,
            singleEvents=True, orderBy='startTime',
            # q='brand.vmeet@nobroker.in' # This might filter too early if brandvmeet is added as resource
        ).execute()
        events = events_result.get('items', [])
        return events
    except HttpError as error:
        print(f'An error occurred fetching events: {error}')
        return []

def load_processed_event_ids():
    # ... (same as before) ...
    if not os.path.exists(PROCESSED_EVENTS_FILE): return set()
    with open(PROCESSED_EVENTS_FILE, 'r') as f: return set(line.strip() for line in f)

def save_processed_event_id(event_id):
    # ... (same as before) ...
    with open(PROCESSED_EVENTS_FILE, 'a') as f: f.write(event_id + '\n')

EVENT_TAG_PROCESSED = "[NBH_BRIEF_AGENT_PROCESSED_V1]"

def is_event_already_tagged(event_description):
    return EVENT_TAG_PROCESSED in (event_description or "")

def tag_event_as_processed(calendar_service, event_id, calendar_id='primary'):
    if not calendar_service:
        print("  Calendar service not available to tag event.")
        return
    try:
        event = calendar_service.events().get(calendarId=calendar_id, eventId=event_id).execute()
        description = event.get('description', '')
        if not is_event_already_tagged(description):
            new_description = f"{description}\n\n{EVENT_TAG_PROCESSED}"
            updated_event_body = {'description': new_description}
            calendar_service.events().patch(
                calendarId=calendar_id, eventId=event_id, body=updated_event_body
            ).execute()
            print(f"  Tagged event {event_id} as processed in calendar.")
    except HttpError as error:
        print(f"  An error occurred tagging event {event_id}: {error}")

EVENT_REMINDER_SET_TAG = "[NBH_1HR_REMINDER_SET]" # New tag

def set_one_hour_email_reminder(calendar_service, event_id, calendar_id='primary'):
    if not calendar_service:
        print("  Calendar service not available to set reminder.")
        return False
    try:
        event = calendar_service.events().get(calendarId=calendar_id, eventId=event_id).execute()
        
        # Check if our specific reminder tag is already there
        description = event.get('description', '')
        if EVENT_REMINDER_SET_TAG in description:
            print(f"  1-hour email reminder already marked as set for event {event_id}.")
            return True

        reminders = event.get('reminders', {})
        overrides = reminders.get('overrides', [])
        
        # Check if a 60-minute email reminder already exists
        has_one_hour_email_reminder = any(
            r.get('method') == 'email' and r.get('minutes') == 60 for r in overrides
        )

        if has_one_hour_email_reminder:
            print(f"  Event {event_id} already has a 60-minute email reminder.")
        else:
            print(f"  Adding 60-minute email reminder to event {event_id}.")
            overrides.append({'method': 'email', 'minutes': 60})
            body_update = {'reminders': {'useDefault': False, 'overrides': overrides}}
            calendar_service.events().patch(
                calendarId=calendar_id, eventId=event_id, body=body_update
            ).execute()
            print(f"  Successfully added 60-minute email reminder for event {event_id}.")

        # Add our custom tag to the description to indicate we've processed this for reminders
        if EVENT_REMINDER_SET_TAG not in description:
            new_description = f"{description}\n{EVENT_REMINDER_SET_TAG}".strip()
            body_update_desc = {'description': new_description}
            calendar_service.events().patch(
                calendarId=calendar_id, eventId=event_id, body=body_update_desc
            ).execute()
            print(f"  Tagged event {event_id} with {EVENT_REMINDER_SET_TAG}.")
        
        return True

    except HttpError as error:
        print(f"  An error occurred setting reminder for event {event_id}: {error}")
        return False



# --- Meeting Info Extraction ---

def extract_meeting_info(event):
    event_id = event['id']
    summary = event.get('summary', 'No Title')
    start_str = event['start'].get('dateTime', event['start'].get('date'))
    start_time_obj = datetime.datetime.fromisoformat(start_str.replace('Z', '+00:00'))
    location = event.get('location', 'N/A')
    description = event.get('description', '')
    attendees = event.get('attendees', [])
    
    nbh_attendees = []
    brand_attendees_info = []
    brand_name_candidates_from_domain = set()

    organizer_email = event.get('organizer', {}).get('email', '').lower()
    if organizer_email and '@nobroker.in' not in organizer_email:
        try:
            brand_name_candidates_from_domain.add(organizer_email.split('@')[1].split('.')[0].capitalize())
        except IndexError: pass

    is_agent_email_accepted = False
    for attendee in attendees:
        email = attendee.get('email','').lower()
        name = attendee.get('displayName', email) 
        if email == AGENT_EMAIL.lower():
            is_agent_email_accepted = True
        elif '@nobroker.in' in email:
            nbh_attendees.append({'email': email, 'name': name})
        elif email:
            brand_attendees_info.append({'name': name, 'email': email})
            if '@' in email:
                try:
                    domain_part = email.split('@')[1]
                    if domain_part.lower() not in ['gmail.com', 'outlook.com', 'yahoo.com', 'hotmail.com', 'aol.com']:
                        brand_name_candidates_from_domain.add(domain_part.split('.')[0].capitalize())
                except IndexError: pass
    
    if not is_agent_email_accepted:
        print(f"  Skipping event '{summary}': {AGENT_EMAIL} has not accepted or is not an attendee.")
        return None

    if not brand_attendees_info:
        return "NO_EXTERNAL_ATTENDEES" 

    brand_name = "Unknown Brand"
    title_lower = summary.lower()
    
    # --- Start Brand Name Extraction ---
    # 1. Define NBH entities and phrases to remove
    nbh_entities = ["nobrokerhood", "nobroker", "nbh","broker","no","hood"]
    meeting_phrases = [
        "meeting with", "call with", "sync with", "internal", "discussion", 
        "intro call", "kick off", "proposal", "review", "update", "connect", 
        "broker", "hood", "team", "session", "chat", "catch up", "briefing", 
        "brief", "sync", "meeting", "call"
    ]
    phrases_to_remove = nbh_entities + meeting_phrases

    # 2. Attempt to specifically handle "BRAND X NBH_ENTITY" or "NBH_ENTITY X BRAND"
    #    This should generally happen on a cleaner version of the title.
    #    Let's first remove common meeting phrases that might interfere.
    
    temp_title_for_x_check = title_lower
    for phrase in meeting_phrases: # Remove only meeting phrases first
        temp_title_for_x_check = temp_title_for_x_check.replace(phrase, " ") # Replace with space

    # Consolidate multiple spaces
    temp_title_for_x_check = re.sub(r'\s+', ' ', temp_title_for_x_check).strip()
    
    extracted_from_x_pattern = False
    for nbh_ent in nbh_entities:
        # Pattern: "brand x nbh_ent" (allow optional spaces around x)
        match = re.search(rf'^(.*?)\s*x\s*{re.escape(nbh_ent)}', temp_title_for_x_check, re.IGNORECASE)
        if match and match.group(1).strip():
            candidate = match.group(1).strip()
            # Clean this candidate further from any remaining NBH entities or general meeting phrases
            for p_to_remove in phrases_to_remove:
                candidate = candidate.replace(p_to_remove, " ").strip()
            candidate = re.sub(r'\s+', ' ', candidate).strip() # Consolidate spaces
            if candidate:
                brand_name = candidate.title()
                extracted_from_x_pattern = True
                break
        
        # Pattern: "nbh_ent x brand"
        match = re.search(rf'^{re.escape(nbh_ent)}\s*x\s*(.*)', temp_title_for_x_check, re.IGNORECASE)
        if match and match.group(1).strip():
            candidate = match.group(1).strip()
            for p_to_remove in phrases_to_remove:
                candidate = candidate.replace(p_to_remove, " ").strip()
            candidate = re.sub(r'\s+', ' ', candidate).strip()
            if candidate:
                brand_name = candidate.title()
                extracted_from_x_pattern = True
                break
        if extracted_from_x_pattern:
            break

    # 3. If not extracted by "X" pattern, proceed with general cleaning and splitting
    current_candidate = title_lower # Start fresh for general processing
    if not extracted_from_x_pattern:
        # Remove all defined phrases (NBH entities and meeting phrases)
        for phrase in phrases_to_remove:
            current_candidate = current_candidate.replace(phrase, " ") # Replace with space to maintain separation

        # Consolidate multiple spaces that might have resulted from replacements
        current_candidate = re.sub(r'\s+', ' ', current_candidate).strip()

        # Define a more comprehensive list of separators.
        # Order can matter. Regex `|` is OR. `\s*` allows zero or more spaces around separators.
        # `re.escape` is important for special characters in separators like `.`
        # Added `<>` and `_` to separators.
        separator_pattern = r'\s*(?:/|-|ft\.|with|and|&|x|:|<>|_)\s*'
        
        potential_parts = re.split(separator_pattern, current_candidate, flags=re.IGNORECASE)
        
        # Filter out empty strings and very short parts after splitting
        valid_parts = [part.strip() for part in potential_parts if part.strip() and len(part.strip()) > 1] # Min length 2 for a part

        if valid_parts:
            # Often, the first part is the brand name if separators are used correctly.
            # However, we also need to consider if a part itself IS a meeting phrase after splitting.
            for part in valid_parts:
                # Check if this part itself is one of the phrases we want to ignore (even after splitting)
                is_a_meeting_phrase = False
                for mp in meeting_phrases: # Check only against meeting_phrases here
                    if part.lower() == mp.lower():
                        is_a_meeting_phrase = True
                        break
                if not is_a_meeting_phrase and len(part) > 1 : # Re-check length
                     # Remove any trailing non-alphanumeric from the chosen part (e.g. trailing month name if it got stuck)
                    clean_part = re.sub(r'[^a-zA-Z0-9\s]+$', '', part).strip() # Remove trailing non-alphanum, non-space
                    # Check if the part is a month name or common day/time terms (to avoid "June" as brand)
                    month_names = ["january", "february", "march", "april", "may", "june", "july", "august", "september", "october", "november", "december",
                                   "jan", "feb", "mar", "apr", "jun", "jul", "aug", "sep", "oct", "nov", "dec", "pm", "am"]
                    if clean_part.lower() not in month_names and not clean_part.isdigit():
                        brand_name = clean_part.title()
                        break 
        elif current_candidate and len(current_candidate) > 1 and not current_candidate.isdigit(): # No separators found, use the whole cleaned candidate
            brand_name = current_candidate.title()

    # 4. Final checks and fallback to domain names
    is_brand_name_generic_or_unknown = (
        brand_name == "Unknown Brand" or
        len(brand_name) <= 1 or # Allow single letter brands if they pass other checks
        brand_name.lower() in phrases_to_remove # Check against full list again
    )

    if is_brand_name_generic_or_unknown and brand_name_candidates_from_domain:
        # Check if domain candidates are not too generic
        filtered_domain_candidates = [
            domain.title() for domain in brand_name_candidates_from_domain 
            if domain.lower() not in phrases_to_remove and len(domain) > 1
        ]
        if filtered_domain_candidates:
            brand_name = " / ".join(sorted(filtered_domain_candidates))
        elif brand_name_candidates_from_domain: # Fallback to original if all were filtered
            brand_name = " / ".join(sorted(list(brand_name_candidates_from_domain)))


    # Final validation of extracted brand_name
    if brand_name == "Unknown Brand" or \
       len(brand_name) <= 1 and brand_name.lower() not in ['x'] or \
       (len(brand_name) > 1 and brand_name.lower() in phrases_to_remove): # Check if it's purely a phrase to remove
        return "AMBIGUOUS_TITLE"

    return {
        'id': event_id, 'title': summary, 'start_time_obj': start_time_obj,
        'start_time_str': start_time_obj.strftime("%Y-%m-%d %I:%M %p %Z (%A)"),
        'location': location, 'description': description, 'nbh_attendees': nbh_attendees,
        'brand_attendees_info': brand_attendees_info, 'brand_name': brand_name.strip().title(), # Ensure title case and stripped
        'is_event_description_present_for_tagging': bool(event.get('description'))
    }


# --- Gemini LLM Integration ---
def configure_gemini():
    if not GEMINI_API_KEY:
        print("GEMINI_API_KEY environment variable not set. LLM will not function.")
        return None
    try:
        genai.configure(api_key=GEMINI_API_KEY)
        # Using a specific model version. 1.5 Flash is faster and cheaper for many tasks.
        # For higher quality, consider 'gemini-1.5-pro-latest'.
        model = genai.GenerativeModel('gemini-2.5-flash-preview-05-20')
        print(f"Gemini model '{model.model_name}' configured successfully.")
        return model
    except Exception as e:
        print(f"Error configuring Gemini API: {e}")
        return None

YOUR_DETAILED_PROMPT_TEMPLATE_GEMINI = """
# NBH is a SAAS platform that digitizes all types of workflows within gated communities such as visitor entry approvals, helpdesk tickets, maintenance payments, communication and announcements etc.. We have 3 platforms, a guard app (used by guards at the gate), a resident app (used by residents in gated communities) and a admin dashboard. We operate in tier-1 cities in India, which means we are engaged with a customer segment which is very valuable to many brands.
# We approach various companies who want to target that segment of people we have for the ads. This might be from a large national/multinational company or a small local city brand.
# You are a world-class marketing strategist at NoBrokerHood (NBH) crafting a Pre-Meeting Brief for an upcoming client pitch for the monetization team. This function has sales team members called brand managers (BMs) that pitch digital assets (on our resident app) and physical monetization assets (society common areas that we monetize) that we have. 
# The goal is to prepare a comprehensive, data-driven brief that will help the BMs present a compelling case to the client.    
Pre-Meeting Brief
# Brand Name: {BRAND_NAME_FOR_BODY}
# Meeting Date and time: {MEETING_DATETIME}
# Brand Attendees:
{BRAND_ATTENDEES_NAMES}
# Attendee from NoBrokerHood (NBH):
{NBH_ATTENDEES_NAMES}

# IMPORTANT: For every section, use markdown headers (##), sub-headers (###), and bullet/numbered lists. Never use asterisks (*) for bullets—use markdown dashes (-) or numbers.

#The full pre meeting brief is divided into two parts.

#Concise executive summary and #Detail Report

# Please format each step as a markdown sub-header (e.g., `## Icebreaking (5 mins)`) followed by a short description or bullet points.
# The total full meeting brief should not cross 6-7 pages. Remember to strictly follow the page count or line count instructions.

# --- CRITICAL INSTRUCTION BASED ON PREVIOUS MEETING DATA ---
# The "Usage of Nobroker/nobrokerhood data" section below will indicate if this meeting with {BRAND_NAME_FOR_BODY}
# appears to be a DIRECT FOLLOW-UP to recent NBH interactions, if NBH has had OTHER PREVIOUS INTERACTIONS
# with the brand, or if this is likely a FIRST-TIME INTERACTION.
# YOUR RESPONSE MUST ADAPT BASED ON THIS:
#   - IF A DIRECT FOLLOW-UP IS INDICATED:
#     - Part 1 (Executive Summary): Keep "Brand & Launch Synopsis" and "Client Profiles" VERY concise, acknowledging the ongoing discussion.
#     - Part 1 (Suggested Meeting Flow): MUST heavily reference previous discussions, focusing on action items or next steps from past meetings.
#     - Part 2 (Detailed Report - Brand section): Can be shorter, focusing on any new developments since the last meeting.
#     - Part 2 (Proposed solution and offering): MUST build upon previously discussed needs or proposals.
#     - AVOID re-introducing NBH basics as if it's the first meeting.
#   - IF OTHER PREVIOUS INTERACTIONS (NOT a direct follow-up) ARE NOTED:
#     - Use insights from past meetings (Key Questions, Brand Traits, Needs, Pain Points) to better inform ALL sections of the current brief,
#       showing NBH's prior understanding. The meeting can still have introductory elements but should be more targeted.
#   - IF NO PREVIOUS INTERACTIONS (First-time meeting):
#     - Generate a full, comprehensive brief as if this is the first major engagement.


Part 1: Concise Executive Summary (keep the concise executive summary strictly below 250 words)

Brand Name: {BRAND_NAME_FOR_BODY}

#Brand & Launch Synopsis: name, launch date/event, one-sentence description of the brand. [IF NOT A DIRECT FOLLOW-UP: name, launch date/event, one-sentence description. IF DIRECT FOLLOW-UP: Briefly acknowledge brand and move on, e.g., "Continuing discussions with Activ8 Pilates X..."]

#Client Profiles (concise - 3 bullet points): Briefly describe their market position, key initiatives, or focus areas. [IF NOT A DIRECT FOLLOW-UP: Describe market position, key initiatives. IF DIRECT FOLLOW-UP: Very briefly summarize current understanding or state "As previously discussed..."]

## Example bullet 1: [e.g., ## Legacy Innovator: Description]
## Example bullet 2: [e.g., ## Market Driver: Description]
## Example bullet 3: [e.g., ## Growth Focus: Description]

#Recent Marketing History (3 bullet points): List 3-4 significant recent marketing campaigns or activities.

## Example bullet 1: [e.g., ## Campaign A: Brief description]
## Example bullet 2: [e.g., ## Campaign B: Brief description]

#Marketing spend bifurcation: completely search the companies financials and the internet and give me the marketing spend bifurcation between traditional media and the new digital media. If the companies bifurcation is not available, give me the industry level bifurcation. Specifically check if there is any spend on activities linked to gated communities or tier 1 cities in India.



#Suggested Meeting Flow (bullet points with timings): (#keep the total time below 30 min). Please include the proposed solution section for sure.

Please format each step as a markdown sub-header (e.g., `## Icebreaking (5 mins)`) followed by a short description or bullet points.

# [CRITICAL: IF A DIRECT FOLLOW-UP, this flow MUST focus on previous action items, unresolved questions, and next steps from the provided "Action Items (Relevant for Follow-up)" data.
# IF NOT A FOLLOW-UP, use the example flow but tailor opening and need-finding based on any prior knowledge.]

Example:

    ## Icebreaking (5 mins)
    - "Good evening, ..."

    ## Need Finding (7 mins)
    - "We understand your platform ..."

    ## Deck Presentation (8 mins)
    - Present NBH's digital and physical assets ...

    #Proposed solution and offering: [IF A DIRECT FOLLOW-UP: Base this heavily on previous discussions and identified needs/action items. Briefly describe creative co-branded solutions or offerings relevant to the client and NoBrokerHood. Try to give atleast 5 solutions. Try to check on the web on some creative marketing campaigns which can be achieved through our collaboration(keep in mind our offerings while giving suggestions.), Keep the points precise  and concise as if you are presenting to a CXO. Keep each description to a maximum of 2 lines. Also check the 3rd slide in the canva ppt. Those are the offerings in the parent company Nobroker.com. We can offer the data in those various verticals as solutioning  to the company and leverage it to get them on board to Nobrokerhood.]
    
    Example Solution 1: [Brief Title/Concept: Description]
    Example Solution 2: [Brief Title/Concept: Description]

    ## Feedback & Brief Taking (3 mins)
    - "Based on what we've shared ..."

    ## Fixing Next Steps (1 min)
    - "Would you be open to ..."

    ## Closing (1 min)
    - "Thank you for your time ..."

Part 2: Detailed report: (keep the detailed report strictly below 600 words.)

#Brand: [IF NOT A DIRECT FOLLOW-UP: Small detailed report on the brand. IF DIRECT FOLLOW-UP: Brief update on any new brand developments or skip if no new info.]

a) categories that the brand exist in

b) Target audience of the brand

c) Peak seasons for the brand where each category sales happen

d) Brand persona

e) What the brand is doing these days and the areas where they are focusing more lately


# Usage of Nobroker/nobrokerhood data: 
    # The following is a summary of relevant internal data from NoBrokerHood's Google Drive. Use this information to identify relevant data points for {BRAND_NAME_FOR_BODY}.

    # This includes:
    #   - Insights from PREVIOUS NBH MEETINGS with {BRAND_NAME_FOR_BODY} (if any). Pay close attention to:
    #       - If the notes indicate a "DIRECT FOLLOW-UP".
    #       - "Key Discussion Points" from past meetings.
    #       - "Action Items (Relevant for Follow-up)" - CRITICAL if this is a follow-up.
    #       - "Key Questions Raised by Client" in past meetings.
    #       - "Observed Brand Traits", "Identified Customer Needs", "Client Pain Points Discussed".
    #     Synthesize these previous meeting insights to inform the entire brief, especially the meeting flow and proposed solutions.    
    #   - Context on our General Pitch Deck.
    #   - Text from the Case Study PDF (analyze to identify distinct case studies: Brand, Objective, Activities, Results, Relevance to {BRAND_NAME_FOR_BODY}), need to also remember the slide numbers for these
    #   - Historical Campaign data for {BRAND_NAME_FOR_BODY} or similar brands/industries.
    #   - A description of available NoBroker.com platform metrics (e.g., leads, user data).
    #
    # BEGIN INTERNAL NBH DATA SUMMARY
    {INTERNAL_NBH_DATA_SUMMARY}
    # END INTERNAL NBH DATA SUMMARY
    # Your task is to synthesize this information:
    #
    # 0. Previous Meeting Insights (If data provided under "Insights from Previous NBH Meetings"):
    #    a. Determine if the current meeting is a direct follow-up based on the notes.
    #    b. If it is a direct follow-up, highlight any "Action Items (Relevant for Follow-up)" that need to be addressed.
    #    c. Summarize any "Key Questions Raised by Client", "Observed Brand Traits", "Identified Customer Needs", or "Client Pain Points Discussed"
    #       from past meetings that are relevant for the upcoming discussion.
    #    d. This understanding should PERVADE the rest of your analysis and suggestions.
    #    e. If no previous meetings are noted, assume this is a first-time interaction and proceed accordingly.
    # 1. Case Studies (from section "Relevant Case Study Slides"):
    #    For EACH "Case Study Slide X" provided:
    #    a. Identify the **Case Study Brand Name** (often in the title or a "Brand:" field).
    #    b. Extract **Campaign Introduction/Objective/Description** (look for "Introduction:", "Description:", or similar phrasing).
    #    c. Extract **Activities/Type of Activity** (look for "Activities:", "Type of Activity:").
    #    d. Extract **Duration** (look for "Duration:").
    #    e. Extract **Key Quantitative Results**. Look for labels like:
    #       - "Revenue:"
    #       - "Leads Generated:" (and "Total Orders" if present)
    #       - "Return on Investment (ROI):"
    #       - "CTR (Click-Through Rate):"
    #       - "Impressions:"
    #       - "Reach:"
    #       - "Footfall:"
    #       (If a specific metric label isn't present but a clear number related to results is, extract it with context.)
    #    f. Extract **City/Location** (look for "City:", "Location:").
    #    g. Based on the extracted information and the current meeting brand ({BRAND_NAME_FOR_BODY}, Industry: [Inferred industry of BRAND_NAME_FOR_BODY if available]), briefly state **why this case study is relevant** (e.g., "Relevant because it's for the same brand," or "Relevant due to similar industry: [Case Study Brand's Industry]," or "Demonstrates success with [specific activity type] which could interest {BRAND_NAME_FOR_BODY}").
    #    h. Also share the case study slide number if the Brand Manager wants to refer it for creatives and other details.
    #   Present the information for each extracted case study clearly, perhaps in a structured format or bullet points under a heading like "Case Study: [Case Study Brand Name] (from Slide X)".
    #    If multiple relevant case studies are provided, select and detail the top 1-3 most impactful or directly relevant ones for {BRAND_NAME_FOR_BODY}. If only one is highly relevant, focus on that.
    #
    # 2. Historical Campaigns (from section "Historical Campaign Data"):
    #    (highlight 1-2 key examples, mention type, duration, results if available for {BRAND_NAME_FOR_BODY} or also mention campaigns done by competitors in the same industry)
    #
    # 3. NoBroker.com Platform Metrics (from section "NoBroker.com Platform Metrics"):
    #    ( You have access to the nobroker data like packers and movers as well as nobrokerhood data like number of maids, pets data etc. Provide the relevant data which we can use to present for this brand {BRAND_NAME_FOR_BODY}. Eg: if its a pet brand, total pet count can be provided, pet count in various cities can be provided to show where the targetting can be done etc. If it is a brand which can be purchased when people move accomodations or cities, move-in move-out data can be cited)
    #
    # Present these insights concisely as actionable data points for the sales team. Ensure to reference the slide number when discussing a specific case study.
    # If information for a specific field (e.g., ROI) is not explicitly found on a slide, state "Not specified on slide." Do not invent data.

#Recent Marketing history: List all the recent marketing history of the brand be it be campaigns or activities. Analyse the recent marketing campaigns of the brand as if you are a McKinsey or a BCG consultant and give me the analysis on various points, like the cohorting they are focusing on, etc. Don’t make a big report. Keep the points precise  and concise as if you are presenting to a CXO. Keep each point to a max of 2 lines.

#Personal History and Persona Analysis: The personal history and persona analysis should be done for all the attendees from the brand side.

## Personal History: You have the name of the brand as well as name of the attendee from the brand. The attendee of the brand is most likely someone from the marketing department of the brand and he can be a product manager, brand manager etc. Having these in mind, conduct a targeted search and you most likely would find the linkedin profile of the attendee. Now, give the output on the basis of the linkedin profile that is obtained. The output should be a maximum of 4 lines.

## Persona Analysis: Using the information we have till now, try to search linkedin, fb, Instagram, twitter and all the social media regarding the person and try to build a persona of the person which contains his hobbies interests, etc. ***Important note: Always give information after confirmation that the information belongs to him/her. And the output should be of maximum 3-4 lines.

Use crisp, professional language and bullet-point formatting. Cite metrics and dates precisely. Produce both the concise summary and the detailed report in one output, clearly labelled.

"""



def generate_brief_with_gemini(gemini_llm_model, meeting_data, internal_data_summary_str):
    if not gemini_llm_model:
        return "Error: Gemini model not available."

    # ... (Format placeholders as in previous thought: nbh_attendee_names_str, brand_attendees_details_str etc.)
    nbh_attendee_names_str = ", ".join([att['name'] for att in meeting_data['nbh_attendees']])
    brand_attendees_info_str = "; ".join([f"{att['name']} ({att['email']})" for att in meeting_data['brand_attendees_info']])
    brand_attendee_names_only_str = ", ".join([att['name'] for att in meeting_data['brand_attendees_info']])


    prompt_filled = YOUR_DETAILED_PROMPT_TEMPLATE_GEMINI.format(
        MEETING_DATETIME=meeting_data['start_time_str'],
        MEETING_LOCATION=meeting_data['location'],
        BRAND_NAME=meeting_data['brand_name'], # This is the overall brand name
        BRAND_ATTENDEES_NAMES=brand_attendee_names_only_str,
        NBH_ATTENDEES_NAMES=nbh_attendee_names_str,
        BRAND_NAME_FOR_BODY=meeting_data['brand_name'], # Used within the prompt body
        BRAND_ATTENDEES_FULL_DETAILS=brand_attendees_info_str,
        INTERNAL_NBH_DATA_SUMMARY=internal_data_summary_str
    )
    
    generation_config = {"temperature": 0.7, "top_p": 0.95, "top_k": 40} # Typical settings
    safety_settings = [
        {"category": "HARM_CATEGORY_HARASSMENT", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
        # ... other safety settings from previous thought block ...
        {"category": "HARM_CATEGORY_HATE_SPEECH", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
        {"category": "HARM_CATEGORY_SEXUALLY_EXPLICIT", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
        {"category": "HARM_CATEGORY_DANGEROUS_CONTENT", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
    ]

    print(f"  Sending request to Gemini for brand: {meeting_data['brand_name']}...")
    try:
        # Use Google Search grounding by adding tools=[{"tool": "google_search"}]
        response = gemini_llm_model.generate_content(
            prompt_filled,
            generation_config=generation_config,
            safety_settings=safety_settings,
            # tools=[{"tool": "google_search"}]  # Enable Google Search grounding >> not working on API calls rights now
        )
        # print(f"Gemini raw response: {response}") # For debugging
        if response.prompt_feedback and response.prompt_feedback.block_reason:
            return f"Error: Prompt blocked by Gemini. Reason: {response.prompt_feedback.block_reason_message or response.prompt_feedback.block_reason}"

        if not response.candidates:
             return f"Error: Gemini returned no candidates. Feedback: {response.prompt_feedback}"

        # Ensure text is extracted correctly
        brief_content = ""
        for part in response.candidates[0].content.parts:
            brief_content += part.text
        
        if not brief_content.strip():
            return "Error: Gemini returned an empty brief."
            
        return brief_content
    except Exception as e:
        print(f"  Error during Gemini API call: {e}")
        # Check for specific API errors if needed, e.g. google.api_core.exceptions.ResourceExhausted
        if "429" in str(e) or "ResourceExhausted" in str(e): # Quota issue
             return "Error: Gemini API quota likely exceeded. Please check your Google Cloud/AI Studio quotas."
        return f"Error: Exception during Gemini call: {e}"

# --- Email Sending ---
def create_email_message(sender, to_emails_list, subject, message_text_html):
    # ... (Using MIMEText with 'html' for better formatting) ...
    message = MIMEText(message_text_html, 'html', 'utf-8') # Send as HTML
    message['to'] = ", ".join(to_emails_list)
    message['from'] = sender
    message['subject'] = subject
    raw_message = base64.urlsafe_b64encode(message.as_bytes())
    return {'raw': raw_message.decode()}

def send_gmail_message(gmail_service, user_id, message_body):
    # ... (same as before) ...
    if not gmail_service:
        print("  Gmail service not available. Cannot send email.")
        return None
    try:
        message = (gmail_service.users().messages().send(userId=user_id, body=message_body).execute())
        print(f'  Message Id: {message["id"]} sent.')
        return message
    except HttpError as error:
        print(f'  An error occurred sending email: {error}')
        return None

def send_brief_email(gmail_service, meeting_data, brief_content):
    nbh_recipient_emails = [att['email'] for att in meeting_data['nbh_attendees'] if att['email'] != AGENT_EMAIL]
    
    # For testing, override recipients:
    nbh_recipient_emails = [ADMIN_EMAIL_FOR_NOTIFICATIONS]
    # print(f"DEBUG: Intended brief recipients: {nbh_recipient_emails}")

    if not nbh_recipient_emails:
        print(f"  No NBH recipients (other than brandvmeet) for '{meeting_data['title']}'. Brief not emailed.")
        return

    email_subject = f"Pre-Meeting Brief: {meeting_data['title']} with {meeting_data['brand_name']}"
    
    # Convert markdown-like brief (from LLM) to basic HTML for email
    html_brief_content = markdown.markdown(brief_content)
    # html_brief_content = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', html_brief_content) # Basic bold
    # Add more markdown to HTML conversions if needed (e.g., for lists, headers)  
    # For headers like # Detail Report, ## Sub-header
    # html_brief_content = html_brief_content.replace("# ", "<h3>").replace("\n", "</h3><br>") # Simple h3 for main headers
    
    email_body_html = f"""
    <html><body>
    <p>Hello Team,</p>
    <p>Please find the pre-meeting brief for your upcoming meeting:</p>
    <p><b>Event:</b> {meeting_data['title']}<br>
    <b>Brand:</b> {meeting_data['brand_name']}<br>
    <b>Scheduled:</b> {meeting_data['start_time_str']}</p>
    <hr>
    {html_brief_content}
    <hr>
    <p>Best regards,<br>NBH Meeting Prep Agent</p>
    </body></html>
    """
    email_message = create_email_message(
        sender=AGENT_EMAIL,
        to_emails_list=nbh_recipient_emails,
        subject=email_subject,
        message_text_html=email_body_html
    )
    send_gmail_message(gmail_service, 'me', email_message)

def send_notification_email(gmail_service, subject, body_html, recipient=ADMIN_EMAIL_FOR_NOTIFICATIONS):
    if not recipient:
        print("  Admin notification email not set. Skipping notification.")
        return
    
    # Always send a copy to brandvmeet for record-keeping
    recipients = list(set([recipient, AGENT_EMAIL]))

    email_message = create_email_message(
        sender=AGENT_EMAIL,
        to_emails_list=recipients,
        subject=subject,
        message_text_html=body_html
    )
    send_gmail_message(gmail_service, 'me', email_message)


# --- Main Execution Logic ---
def main():
    print(f"Script started at {datetime.datetime.now()}")
    print(f"Using NBH GDrive Folder ID: {NBH_GDRIVE_FOLDER_ID}")

    # Initialize Google Services
    # Use a combined token file strategy or separate ones. Separate is fine.
    calendar_service = get_google_service('calendar', 'v3', SCOPES, f"{TOKEN_FILE_PREFIX}_calendar.json")
    gmail_service = get_google_service('gmail', 'v1', SCOPES, f"{TOKEN_FILE_PREFIX}_gmail.json")
    drive_service = get_google_service('drive', 'v3', SCOPES, f"{TOKEN_FILE_PREFIX}_drive.json")
    sheets_service = get_google_service('sheets', 'v4', SCOPES, f"{TOKEN_FILE_PREFIX}_sheets.json")
    gemini_llm_model = configure_gemini()


    # --- REMOVED: Pre-processing of campaign sheets for all brand industries ---
    # --- Pre-process campaign sheets to get all brand industries (uses cache) ---
   # campaign_sheets_to_process_for_industries = [
    #    {'file_name_keyword': FILE_NAME_PHYSICAL_CAMPAIGNS_GSHEET.lower(), 'brand_name_column_header': "Brand Name"},
    #    {'file_name_keyword': FILE_NAME_DIGITAL_CAMPAIGNS_GSHEET.lower(), 'brand_name_column_header': "Brand Name"}
    #]


    # Pass "" for sheets_service if it's already initialized above
    #all_campaign_brands_industry_map = {}
    #if drive_service and sheets_service: # LLM availability checked inside the function
    #    all_campaign_brands_industry_map = get_and_infer_industries_for_all_campaign_brands(
    #        drive_service,
    #        sheets_service,
    #        gemini_llm_model,
    #        campaign_sheets_to_process_for_industries,
    #        batch_size=25 # Optionally override default batch_size here
    #    )
    
    #elif os.path.exists(INFERRED_INDUSTRIES_CACHE_FILE): # If drive/sheets not available, but cache exists
    #    all_campaign_brands_industry_map = load_inferred_industries_cache()
    #    print("Drive/Sheets service not fully available. Loaded industries from existing cache.")
    #else: # No services, no cache
    #    print("Drive/Sheets service not fully available and no cache found. Industry map will be empty.")


    

    if not calendar_service: # Critical service
        print("Exiting: Calendar service failed to initialize.")
        return

    # Load internal data once per run
    #internal_nbh_data_summary = "Internal NBH Data: Not fetched (Drive service issue or not configured)."
    #if drive_service and gemini_llm_model: # Only fetch if we can use it
    #     internal_nbh_data_summary = get_internal_nbh_data_summary(drive_service,sheets_service)
    #elif not gemini_llm_model:
    #    internal_nbh_data_summary = "Internal NBH Data: Not fetched (LLM not available)."


    upcoming_events = get_upcoming_meetings(calendar_service)
    if not upcoming_events:
        print('No upcoming events found for brand.vmeet@nobroker.in that need processing.')
        return

    processed_ids_local_file = load_processed_event_ids()

    for event_payload in upcoming_events:
        event_id = event_payload['id']
        event_summary = event_payload.get('summary', 'No Title')
        event_description_for_tag_check = event_payload.get('description')

        print(f"\nProcessing event: '{event_summary}' (ID: {event_id})")

        # Robust check for "already processed" using calendar event tags
        if is_event_already_tagged(event_description_for_tag_check):
            print(f"  Skipping event '{event_summary}': Already tagged as processed in calendar description.")
            continue
        
        # Fallback for local runs if tagging isn't working or not yet implemented fully
        if event_id in processed_ids_local_file:
            print(f"  Skipping event '{event_summary}': Found in local processed_event_ids.txt (might be redundant if tagging works).")
            continue

        meeting_data_result = extract_meeting_info(event_payload)

        if meeting_data_result is None: # brandvmeet not accepted
            save_processed_event_id(event_id) # Mark locally to avoid re-evaluating simple skips
            tag_event_as_processed(calendar_service, event_id) # Also tag in calendar
            continue 
        
        if meeting_data_result == "NO_EXTERNAL_ATTENDEES":
            print(f"  Event '{event_summary}': No external attendees. No brief needed.")
            save_processed_event_id(event_id)
            tag_event_as_processed(calendar_service, event_id)
            continue
        
        if meeting_data_result == "AMBIGUOUS_TITLE":
            print(f"  Event '{event_summary}': Title is ambiguous for brand extraction.")
            ambiguous_body_html = f"""
            <html><body><p>The pre-meeting brief agent could not reliably determine the brand for the meeting:</p>
            <p><b>Event:</b> {event_summary}<br>
            <b>Scheduled:</b> {event_payload['start'].get('dateTime', event_payload['start'].get('date'))}</p>
            <p>To ensure briefs are generated correctly, please use clear and consistent meeting titles. Suggestions:
            <ul><li>'NBH / [Brand Name] - Kick-off'</li><li>'[Brand Name] & NoBrokerHood - Proposal Discussion'</li></ul></p>
            <p>No brief was generated for this meeting.</p></body></html>"""
            send_notification_email(gmail_service, 
                                    f"Action Required: Ambiguous Title for Meeting - {event_summary}",
                                    ambiguous_body_html)
            save_processed_event_id(event_id)
            tag_event_as_processed(calendar_service, event_id)
            continue

        # If we reached here, meeting_data_result is the actual data dictionary
        meeting_data = meeting_data_result
        current_brand_name_for_meeting = meeting_data['brand_name']

        internal_nbh_data_for_brand = "Internal NBH Data: Not fetched or processed."


        if drive_service and sheets_service: # LLM model is now passed to get_internal_nbh_data_for_brand
            internal_nbh_data_for_brand_str = get_internal_nbh_data_for_brand(
                drive_service, 
                sheets_service, 
                gemini_llm_model, 
                current_brand_name_for_meeting,
                current_meeting_data=meeting_data # PASS THE FULL DICT
            )

        if not gemini_llm_model:
            print(f"  Skipping brief generation for '{meeting_data['title']}': Gemini LLM not available.")
            # Don't mark as processed yet, maybe LLM will be available next run
            continue
        
        if not meeting_data['nbh_attendees']: # Check if any NBH humans are there
            print(f"  Event '{meeting_data['title']}': No NBH attendees (other than brandvmeet) to send brief to.")
            save_processed_event_id(event_id)
            tag_event_as_processed(calendar_service, event_id)
            continue


        print(f"  Proceeding with brief generation for: {meeting_data['brand_name']}")
        generated_brief = generate_brief_with_gemini(gemini_llm_model, meeting_data, internal_nbh_data_for_brand_str)

        if "Error:" in generated_brief: # Check for errors from LLM
            print(f"  Failed to generate brief for '{meeting_data['title']}': {generated_brief}")
            error_body_html = f"""
            <html><body><p>The pre-meeting brief agent encountered an error while generating the brief for:</p>
            <p><b>Event:</b> {meeting_data['title']}<br>
            <b>Brand:</b> {meeting_data['brand_name']}<br>
            <b>Scheduled:</b> {meeting_data['start_time_str']}</p>
            <p><b>Error details:</b> {generated_brief}</p></body></html>"""
            send_notification_email(gmail_service,
                                    f"Error Generating Brief: {meeting_data['title']}",
                                    error_body_html)
            # Don't tag as fully processed if LLM fails, maybe it's temporary.
            # Or use a different tag like [NBH_BRIEF_AGENT_ERROR_V1]
        else:
            print(f"  Successfully generated brief for '{meeting_data['title']}'.")
            send_brief_email(gmail_service, meeting_data, generated_brief)
            tag_event_as_processed(calendar_service, event_id) # Tag only on full success
            set_one_hour_email_reminder(calendar_service, event_id) # Add this call
            save_processed_event_id(event_id)
        

    print(f"Script finished at {datetime.datetime.now()}")

if __name__ == '__main__':
    main()
