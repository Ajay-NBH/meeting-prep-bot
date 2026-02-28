import datetime
import os.path
import time
import base64
import traceback
from email.mime.text import MIMEText
import io # For GDrive downloads
import re
import markdown 
import json
import fitz
from google import genai
from google.genai import types
import enum
import pandas as pd

from google.auth.transport.requests import Request
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from dotenv import load_dotenv
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from googleapiclient.http import MediaIoBaseDownload
from dotenv import load_dotenv
from pydantic import BaseModel, ValidationError
from data_config import sheet_masters, hierarchy, column_index

# For parsing Office documents if downloaded from Drive
from pptx import Presentation
import openpyxl

# Load the .env file
env = os.getenv("ENV", "dev")
env_file = f".env.{env}"

if not os.getenv("GITHUB_ACTIONS"):
    load_dotenv(env_file)

if not os.getenv("GITHUB_ACTIONS"):
    load_dotenv(env_file)
# --- Configuration ---
# For Google Workspace APIs (Calendar, Gmail, Drive)
SCOPES = [
    'https://www.googleapis.com/auth/calendar',
    'https://www.googleapis.com/auth/gmail.send',
    'https://www.googleapis.com/auth/drive.readonly',
    'https://www.googleapis.com/auth/spreadsheets',
    'https://www.googleapis.com/auth/drive',
    'https://www.googleapis.com/auth/documents'
]
CREDENTIALS_FILE = os.getenv("GOOGLE_CREDENTIALS_FILE", "credentials.json") # Downloaded from GCP
TOKEN_FILE_PREFIX = 'token_brandvmeet' # Will generate token_brandvmeet_calendar.json etc.

# For Gemini API
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY") # Set this environment variable

# Google Drive Folder ID containing NBH data
NBH_GDRIVE_FOLDER_ID = os.getenv("NBH_GDRIVE_FOLDER_ID") # Set env var or replace placeholder

AGENT_EMAIL = "brand.vmeet@nobroker.in" # Email of the agent account
ADMIN_EMAIL_FOR_NOTIFICATIONS = "ajay.saini@nobroker.in" # REPLACE with your actual email
leadership_emails = ["sristi.agarwal@nobroker.in", "rohit.c@nobroker.in"] # Add the second email

EXCLUDED_NBH_PSEUDO_NAMES_FOR_FOLLOWUP = {
    AGENT_EMAIL.lower().split('@')[0],
    "pia.brand","pia","nbh.meeting"
}

PROCESSED_EVENTS_FILE = 'processed_event_ids.txt' # Simple file-based tracking for local runs



# Cache file for inferred industries
# INFERRED_INDUSTRIES_CACHE_FILE = 'inferred_industries_cache.json' # Cache for brand industry inference



# Specific File Names (you might make these configurable or discover them)
FILE_NAME_PITCH_DECK_PDF = "NBH Monetization Pitch Deck.pdf" # From your image
# FILE_NAME_CASE_STUDIES_GSLIDES = "National Campaigns_case studies" # From your image
FILE_NAME_CASE_STUDIES_PDF = "National_Campaigns_case_studies.pdf" # New or alternative
FILE_NAME_PHYSICAL_CAMPAIGNS_GSHEET = "Physical_campaigns_live_sheet" # From your image
FILE_NAME_DIGITAL_CAMPAIGNS_GSHEET = "Digital_Campaigns_live_sheet" # From your image (note "live sheet")
FILE_NAME_COM_DATA_GSHEET = "NoBroker_Overall_Data" # From your image
FILE_NAME_NBH_PREVIOUS_MEETINGS_GSHEET = "NBH_previous_meetings_updated"
FILE_NAME_LATEST_CASE_STUDIES_GSHEET = "Consolidated Case Studies - Master"

def parse_names_from_cell_helper(cell_value_str):
    """
    Intelligently parses a string from a spreadsheet cell to extract a set of cleaned, lowercase names.
    Handles two primary formats:
    1. A string representation of a list (e.g., "['name1@example.com', 'name2']").
    2. A simple delimited string (e.g., "Name One, Name Two & Name Three").
    """
    names = set()
    if not cell_value_str or str(cell_value_str).strip().lower() == 'n/a':
        return names

    # Ensure we are working with a string and clean it up
    cell_value_str = str(cell_value_str).strip()
    potential_names = []

    # --- NEW LOGIC: Check if the string is formatted like a list ---
    if cell_value_str.startswith('[') and cell_value_str.endswith(']'):
        # Use regex to find all content within single or double quotes
        # This is safer than using eval()
        extracted_items = re.findall(r"[\'\"](.*?)[\'\"]", cell_value_str)
        for item in extracted_items:
            # If the item is an email, just take the part before the @
            if '@' in item:
                potential_names.append(item.split('@')[0])
            else:
                potential_names.append(item)
    else:
        # --- FALLBACK LOGIC: Handle simple, delimited strings ---
        # Remove content in parentheses (e.g., "(NoBrokerHood)")
        cleaned_cell = re.sub(r'\s*\([^)]*\)', '', cell_value_str)
        cleaned_cell = cleaned_cell.replace('*', '').strip() # Remove asterisks

        # Split by common delimiters
        potential_names = re.split(r'\s*[,;/&\n]\s*|\s+\band\b\s+|\s+\bwith\b\s+', cleaned_cell)

    # --- Common cleaning process for all extracted name parts ---
    for name_part in potential_names:
        final_name = name_part.strip().lower()
        # Filter for meaningful names and exclude common role descriptors
        if final_name and len(final_name) > 2 and \
           "nbh sales" not in final_name and \
           "brand representative" not in final_name and \
           "nobrokerhood" not in final_name and \
           "stay vista" not in final_name:
            names.add(final_name)
            
    return names

# --- ADD THESE NEW HELPER FUNCTIONS NEAR THE TOP OF YOUR SCRIPT ---

def normalize_attendee_name(name_str):
    """
    Takes a raw name string and converts it into a standardized set of name parts.
    - "Shubham Chandrakant Dakhane" -> {'shubham', 'chandrakant', 'dakhane'}
    - "shubham.chandrakant" -> {'shubham', 'chandrakant'}
    - "trisha.bagchi7" -> {'trisha', 'bagchi'}
    - "mary" -> {'mary'}
    """
    if not isinstance(name_str, str) or not name_str.strip():
        return set()
    
    # Lowercase and replace common delimiters with spaces
    processed_name = name_str.lower().replace('.', ' ').replace('_', ' ')
    
    # Remove all characters that are not letters or spaces
    processed_name = re.sub(r'[^a-z\s]', '', processed_name)
    
    # Split into parts and filter out any empty strings resulting from multiple spaces
    name_parts = {part for part in processed_name.split() if part}
    
    return name_parts

def find_common_attendees(attendee_set_1_raw, attendee_set_2_raw):
    """
    Compares two sets of raw name strings and finds common individuals
    using a flexible, normalization-based approach.
    Returns a list of the matched raw names from the first set.
    """
    # Normalize all names in both sets
    # Each item in these lists will be a set of name parts, e.g., [{'shubham', 'chandrakant'}, {'trisha', 'bagchi'}]
    normalized_attendees_1 = [normalize_attendee_name(name) for name in attendee_set_1_raw]
    normalized_attendees_2 = [normalize_attendee_name(name) for name in attendee_set_2_raw]

    common_attendees_raw_names = []
    
    # Keep track of which attendees from set 2 have already been matched to avoid double counting
    matched_indices_in_set_2 = set()

    for i, norm_set_1 in enumerate(normalized_attendees_1):
        if not norm_set_1:
            continue
        
        for j, norm_set_2 in enumerate(normalized_attendees_2):
            if j in matched_indices_in_set_2 or not norm_set_2:
                continue

            # --- NEW, MORE ROBUST CORE LOGIC ---
            # A match occurs if:
            # 1. The name sets are identical (e.g., {'john', 'doe'} == {'john', 'doe'}).
            # 2. One set is a complete subset of the other (e.g., {'john'} is a subset of {'john', 'doe'}).
            # 3. There is a non-trivial intersection (e.g., {'john', 'd'} and {'john', 'doe'} intersect on 'john').
            #    We add the len > 1 check to avoid matching on single initials like 'a'.

            intersection = norm_set_1.intersection(norm_set_2)
            is_match = False
            
            if norm_set_1 == norm_set_2:
                is_match = True
            elif norm_set_1.issubset(norm_set_2) or norm_set_2.issubset(norm_set_1):
                is_match = True
            elif intersection and any(len(name_part) > 1 for name_part in intersection):
                is_match = True

            if is_match:
                common_attendees_raw_names.append(list(attendee_set_1_raw)[i])
                matched_indices_in_set_2.add(j)
                break # Match found, move to the next person in set 1

    return common_attendees_raw_names

# --- END OF NEW HELPER FUNCTIONS ---

# ========== NEW FUNCTION 1: Search for LinkedIn Profile ==========
def search_linkedin_profile(person_name, company_name, gemini_llm_client):
    """
    Uses Gemini's Google Search to find a LinkedIn URL with improved query logic.
    """
    if not gemini_llm_client:
        return None
    
    # Simpler, broader search prompt that allows the LLM to find the best match
    search_prompt = f"Find the official LinkedIn profile URL for {person_name} who works at {company_name} in India. Search for their current role. Return ONLY the URL."

    grounding_tool = types.Tool(google_search=types.GoogleSearch())
    config = types.GenerateContentConfig(
        temperature=0.0,
        tools=[grounding_tool]
    )
    
    try:
        response = gemini_llm_client.models.generate_content(
            model="gemini-2.0-flash", # Use 2.0 for better search grounding
            contents=search_prompt,
            config=config
        )
        
        result_text = response.text.strip()
        
        # Flexible URL pattern matching for all LinkedIn variations
        url_pattern = r'https?://(?:[a-z]+\.)?linkedin\.com/in/[a-zA-Z0-9%_-]+'
        urls = re.findall(url_pattern, result_text)
        
        return urls[0] if urls else None
        
    except Exception as e:
        print(f"  Error searching LinkedIn for {person_name}: {e}")
        return None


# ========== NEW FUNCTION 2: Get LinkedIn for All Attendees ==========
def get_brand_attendees_linkedin_info(brand_attendees_list, brand_name, gemini_llm_client):
    """
    For each brand attendee, search for their LinkedIn profile.
    Returns a list with LinkedIn URLs added.
    """
    attendees_with_linkedin = []
    
    for attendee in brand_attendees_list:
        attendee_name = attendee.get('name', '')
        attendee_email = attendee.get('email', '')
        
        # Clean up name if it looks like an email
        if '@' in attendee_name:
            # Convert "john.doe@example.com" to "John Doe"
            attendee_name = attendee_email.split('@')[0].replace('.', ' ').replace('_', ' ').title()
        
        print(f"    🔍 Searching LinkedIn for: {attendee_name} at {brand_name}")
        
        # Do the actual LinkedIn search
        linkedin_url = search_linkedin_profile(attendee_name, brand_name, gemini_llm_client)
        
        # Add to results
        attendees_with_linkedin.append({
            'name': attendee_name,
            'email': attendee_email,
            'linkedin_url': linkedin_url if linkedin_url else '(LinkedIn Not Verified)'
        })
        
        # Wait 10 seconds to avoid hitting 429 Rate Limits
        print("    ⏳ Waiting 10s to respect API quota...")
        time.sleep(10)
    
    return attendees_with_linkedin

# ========== NEW FUNCTION 3: Find Potential Key Contacts (FIXED) ==========
def find_potential_key_contacts(brand_name, gemini_llm_client):
    """
    Finds 2-3 current Brand Managers or Program Managers at the company in India.
    """
    if not gemini_llm_client:
        return []
    
    discovery_prompt = f"""
Use Google Search to find 2-3 current marketing or brand professionals at {brand_name} India.

Target roles (Priority order):
1. Brand Manager / Senior Brand Manager
2. Marketing Program Manager
3. Media Lead / Digital Marketing Manager

Search Strategy:
- Focus specifically on "{brand_name} India LinkedIn Brand Manager"
- Focus specifically on "{brand_name} India Marketing Program Manager"

IMPORTANT RULES:
1. Identify REAL people currently working in these roles in India.
2. Return in this EXACT JSON format:
{{
  "contacts": [
    {{"name": "Full Name", "title": "Job Title", "reasoning": "Manages brand programs and marketing budgets"}}
  ]
}}
3. If no one is found, return {{"contacts": []}}
"""
    grounding_tool = types.Tool(google_search=types.GoogleSearch())
    config = types.GenerateContentConfig(
        temperature=0.1, # Lower temperature for better accuracy
        tools=[grounding_tool],
        response_mime_type="application/json"
    )
    
    try:
        # Step 1: Discover people using Google Search Grounding
        response = gemini_llm_client.models.generate_content(
            model="gemini-2.0-flash", # Improved grounding model
            contents=discovery_prompt,
            config=config
        )
        
        result_text = response.text.strip()
        result_text = re.sub(r'```json\s*|\s*```', '', result_text).strip()
        contacts_data = json.loads(result_text)
        discovered_contacts = contacts_data.get("contacts", [])
        
        enriched_contacts = []
        for contact in discovered_contacts[:3]: # Max 3 contacts
            name = contact.get("name", "")
            if not name: continue
            
            print(f"      Searching LinkedIn for discovered contact: {name}")
            
            # Step 2: Use search_linkedin_profile to get the specific URL
            linkedin_url = search_linkedin_profile(name, brand_name, gemini_llm_client)
            
            enriched_contacts.append({
                'name': name,
                'title': contact.get("title", ""),
                'reasoning': contact.get("reasoning", "Key decision-maker for NBH collaborations"),
                'linkedin_url': linkedin_url if linkedin_url else '(LinkedIn Not Verified)'
            })
            time.sleep(2) # Pause between searches to avoid hitting quotas
        
        return enriched_contacts
    except Exception as e:
        print(f"    Error in key contact discovery for {brand_name}: {e}")
        return []

class Industry(enum.Enum):
    FMCG = "FMCG"
    AUTOMOTIVE_AND_TRANSPORT = "Automotive & Transportation"
    MEMBERSHIP_AND_LOCAL_SERVICES="Membership & Local Services"
    MARKETING_ADVERTISING_AND_MEDIA="Marketing, Advertising & Media"
    APPAREL_AND_FASHION="Apparel & Fashion"
    FOOD_AND_BEVERAGE="Food & Beverage"
    HEALTHCARE="Healthcare"
    FINANCE_AND_FINTECH="Finance & Fintech"
    BEAUTY_AND_PERSONAL_CARE="Beauty & Personal Care"
    JEWELLERY="Jewellery"
    REAL_ESTATE_AND_CONSTRUCTION="Real Estate & Construction"
    ENERGY_RENEWABLES_AND_MINING="Energy, Renewables & Mining"
    WELLNESS_AND_FITNESS="Wellness & Fitness"
    EDUCATION_AND_TRAINING="Education & Training"
    HOME_GOODS_AND_ELECTRONINCS="Home Goods & Electronics"
    HOSPITALITY_AND_TRAVEL="Hospitality & Travel"
    TECHNOLOGY_AND_BUSINESS_SERVICES="Technology & Business Services"
    E_COMMERCE="E-Commerce"
    RETAIL="Retail"
    PETS_AND_PETS_SERVICES="Pets & Pet Services"
    GAMING="Gaming"
    LOGISTICS_AND_WAREHOUSING="Logistics & Warehousing"
    OTHER_UNKNOWN="Other / Unknown"
    MANUFACTURING_AND_INDUSTRIAL="Manufacturing & Industrial"
    QUICk_COMMERCE="Quick Commerce"
    PHARMA = "Pharma"
    OTT = "OTT"



class Brand_Details(BaseModel):
    brand_name: str
    industry:   Industry
    
    class Config:
        use_enum_values = True  # Use enum values instead of names in JSON output

Allowed_Industries = [industry.value for industry in Industry]

BRAND_EXTRACTION_PROMPT_TEMPLATE = """
You are an expert administrative assistant working for NoBrokerHood (NBH) responsible for parsing meeting titles of meetings between NBH and different companies to extract key business information about those companies.
Your task is to analyze the provided meeting title and return a JSON object with two specific keys: "brand_name" and "industry".

Follow these rules precisely:
1.  **brand_name**: Identify the primary brand or company being met. If a title follows the pattern 'Parent Company (Brand)', the 'Brand' inside the parentheses is the primary `brand_name`. The parent company should be ignored for this task.
2.  **industry**: Infer the most likely industry for the primary `brand_name` strictly from **Allowed_Industries** mentioned below.
3.  If the title is ambiguous or you cannot identify a clear brand, return "Unknown" for both fields.
4.  Your response MUST be ONLY the JSON object, with no other text or markdown fences.

Allowed_Industries: {Allowed_Industries}
---
Here are some examples:

Title: "TCPL (Tetley) X NoBrokerHood/Partnership, 11am"
{{
  "brand_name": "Tetley",
  "industry": "Beverages"
}}

Title: "NBH X GIVA Digital _ June Discussion, 12pm"
{{
  "brand_name": "Giva",
  "industry": "Jewellery"
}}

Title: "Physical meeting - Posterscope X NBH, 3:15pm"
{{
  "brand_name": "Posterscope",
  "industry": "Media Agency"
}}

Title: "Campaign Discussion | Aashirvaad Svasti x NBH, 5pm"
{{
  "brand_name": "Aashirvaad Svasti",
  "industry": "Dairy"
}}

Title: "Internal Team Sync"
{{
  "brand_name": "Unknown",
  "industry": "Unknown"
}}
---

Now, analyze the following title:

Title: "{MEETING_TITLE}"
"""

def get_brand_details_from_title_with_llm(gemini_llm_client, meeting_title):
    """
    Uses a single LLM call to extract brand name and industry from a meeting title.
    Returns a dictionary with the extracted info or defaults if parsing fails.
    """
    default_response = {
        "brand_name": "Unknown Brand",
        "industry": "Unknown"
    }
    if not gemini_llm_client:
        print("  LLM model not available for brand extraction.")
        return default_response

    prompt = BRAND_EXTRACTION_PROMPT_TEMPLATE.format(MEETING_TITLE=meeting_title, Allowed_Industries=Allowed_Industries)

    # Define the grounding tool
    grounding_tool = types.Tool(
        google_search=types.GoogleSearch()
    )

    # Configure generation settings
    config = types.GenerateContentConfig(
        tools=[grounding_tool]
    )
    try:
        raw_text = "" # <-- ADD THIS LINE
        response = gemini_llm_client.models.generate_content(model="gemini-2.5-flash",contents=prompt, config=config)
        raw_text = response.candidates[0].content.parts[0].text

        cleaned_json_str = re.sub(r'```json\s*|\s*```', '', raw_text).strip()
        data = json.loads(cleaned_json_str)

        # Validate the simplified response structure
        if "brand_name" in data and "industry" in data:
            if not data["brand_name"] or data["brand_name"].lower() == 'unknown':
                print(f"  LLM identified title '{meeting_title}' as ambiguous.")
                return default_response
            return data
        else:
            print(f"  Error: LLM response for '{meeting_title}' was missing 'brand_name' or 'industry'.")
            return default_response

    except (json.JSONDecodeError, IndexError, AttributeError, Exception) as e:
        print(f"⚠️ First pass failed ({e}), retrying with strict JSON schema…")

        # Build a cleanup prompt + strict JSON enforcement (no tools)
        cleanup_prompt = (
            "The text below is supposed to be a JSON object matching this schema:\n\n"
            f"{json.dumps(Brand_Details.model_json_schema(), indent=2)}\n\n"
            "But it wasn’t valid JSON. Please reformat exactly as JSON (no extra text):\n\n"
            f"{raw_text}"
        )

        cleanup_config = types.GenerateContentConfig(
            response_mime_type="application/json",
            response_schema=Brand_Details
        )
        retry = gemini_llm_client.models.generate_content(
            model="gemini-2.5-flash",
            contents=cleanup_prompt,
            config=cleanup_config
        )

        # the SDK will give you a .parsed attribute when you supply response_schema
        try:
            parsed: Brand_Details = retry.parsed
            return parsed.model_dump()   # or parsed.dict()
        except Exception as e2:
            print(f"❌ Retry still failed: {e2}")
            return default_response

    except Exception as e:
        # anything else
        print(f"❌ Unexpected error: {e}")
        return default_response


# --- Google Authentication and Service Building ---
def get_google_service(service_name, version, scopes_list, token_filename_base_for_local_storage): # Changed last param name for clarity
    creds = None
    # Construct the specific local token filename (e.g., token_brandvmeet_calendar.json)
    # This is used for local development fallback and saving tokens locally.
    local_token_file_path = token_filename_base_for_local_storage

    # --- Attempt 1: Load from specific environment variable for the service (for CI) ---
    # Construct the expected environment variable name, e.g., GOOGLE_TOKEN_JSON_CALENDAR
    token_env_var_name = f"GOOGLE_TOKEN_JSON_{service_name.upper()}"
    token_json_string_from_env = os.getenv(token_env_var_name)

    if os.getenv('CI') == 'true' and token_json_string_from_env: # Check 'CI' env var and if token string exists
        print(f"CI environment detected. Attempting to load credentials for {service_name} from env var: {token_env_var_name}")
        try:
            token_info = json.loads(token_json_string_from_env)
            # The token_info from your stored JSON should contain client_id, client_secret, and refresh_token,
            # which are needed for the refresh mechanism by the Credentials object.
            creds = Credentials.from_authorized_user_info(token_info, scopes_list)
            print(f"Successfully loaded credentials for {service_name} from environment variable.")
        except json.JSONDecodeError as e:
            print(f"Error decoding JSON from {token_env_var_name} for {service_name}: {e}")
            creds = None
        except Exception as e:
            print(f"Generic error loading credentials for {service_name} from environment variable {token_env_var_name}: {e}")
            creds = None # Fallback

    # --- Attempt 2: Load from local token file (for local development or if CI load failed) ---
    if not creds and os.path.exists(local_token_file_path):
        print(f"Loading credentials for {service_name} from local file: {local_token_file_path}")
        try:
            creds = Credentials.from_authorized_user_file(local_token_file_path, scopes_list)
        except Exception as e:
            print(f"Error loading credentials from {local_token_file_path} for {service_name}: {e}")
            creds = None

    # --- Attempt 3 & 4: Refresh or Run Interactive Flow ---
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            print(f"Token for {service_name} is expired. Refreshing...")
            try:
                creds.refresh(Request())
                print(f"Token for {service_name} refreshed successfully.")
                # If refreshed successfully in a non-CI environment, save it back to the local file
                if not os.getenv('CI') == 'true':
                    try:
                        with open(local_token_file_path, 'w') as token_file:
                            token_file.write(creds.to_json())
                        print(f"Saved refreshed token for {service_name} to {local_token_file_path}")
                    except Exception as e_save:
                        print(f"Error saving refreshed token for {service_name} to {local_token_file_path}: {e_save}")
            except Exception as e_refresh:
                print(f"Error refreshing token for {service_name}: {e_refresh}")
                # If refresh fails in CI, it's a critical issue as there's no interactive fallback.
                if os.getenv('CI') == 'true':
                    print(f"FATAL: Token refresh failed for {service_name} in CI. Stored token might be invalid (e.g., revoked) or scopes changed.")
                    return None
                creds = None # Force re-authentication locally if refresh failed and not in CI
        
        # This block should ideally NOT be reached in a CI environment if the token env var is set correctly.
        if not creds:
            if os.getenv('CI') == 'true':
                print(f"FATAL: No valid credentials for {service_name} in CI and interactive flow is disabled.")
                print(f"       Ensure {token_env_var_name} secret is set correctly and contains valid token JSON.")
                return None # Critical failure in CI

            print(f"No valid credentials for {service_name}. Attempting interactive local server flow...")
            # Ensure CREDENTIALS_FILE ('credentials.json' with OAuth client_id/secret) exists for the flow
            if not os.path.exists(CREDENTIALS_FILE):
                print(f"FATAL: {CREDENTIALS_FILE} not found. Cannot run interactive auth flow for {service_name}.")
                return None
            try:
                flow = InstalledAppFlow.from_client_secrets_file(CREDENTIALS_FILE, scopes_list)
                creds = flow.run_local_server(port=0) # This opens a browser for user consent
                # Save the newly obtained credentials to the local token file for future runs
                with open(local_token_file_path, 'w') as token_file:
                    token_file.write(creds.to_json())
                print(f"Saved new token for {service_name} (from interactive flow) to {local_token_file_path}")
            except Exception as e_flow:
                print(f"Error during interactive auth flow for {service_name}: {e_flow}")
                return None
    
    if not creds: # Should not happen if all paths above are handled
        print(f"Ultimately failed to obtain credentials for {service_name}.")
        return None

    # --- Build the Google API Service ---
    try:
        service = build(service_name, version, credentials=creds)
        print(f"{service_name.capitalize()} service initialized successfully.")
        return service
    except HttpError as error:
        print(f'An HTTP error occurred building {service_name} service: {error}')
        return None
    except Exception as e: # Catch other potential errors during build
        print(f'A general error occurred building {service_name} service: {e}')
        return None
    
# --- Google Drive Functions ---
def list_files_in_gdrive_folder(drive_service, folder_id):
    # ... (Implementation from previous thought block - list files) ...
    # Ensure it handles empty folder_id
    if not folder_id or folder_id == "YOUR_GDRIVE_FOLDER_ID_HERE":
        print("Google Drive Folder ID for NBH data is not configured.")
        return []
    try:
        results = drive_service.files().list(
            q=f"'{folder_id}' in parents and trashed=false",
            # Include mimeType to decide how to read
            fields="nextPageToken, files(id, name, mimeType)"
        ).execute()
        items = results.get('files', [])
        return items
    except HttpError as error:
        print(f"An error occurred listing GDrive files for folder {folder_id}: {error}")
        return []

    

# THIS IS THE PRIMARY FUNCTION TO GET FILE DATA
def get_structured_gdrive_file_data(drive_service, sheets_service, file_id, file_name, mime_type):
    """
    Extracts and parses structured content from a Google Drive file based on its MIME type.
    
    Supports Google Slides (as PPTX), Microsoft PowerPoint, Google Sheets, Microsoft Excel, Google Docs, PDFs, and plain text files. Returns structured data as a list of dictionaries for presentations and spreadsheets, or as a string for text-based files. If parsing fails or the file type is unsupported, returns a descriptive error message.
    """

    print(f"    Attempting to read structured data for: {file_name} (MIME: {mime_type})")

    MIMETYPE_GOOGLE_SHEET = 'application/vnd.google-apps.spreadsheet'
    MIMETYPE_GOOGLE_DOC = 'application/vnd.google-apps.document'
    MIMETYPE_GOOGLE_PRESENTATION = 'application/vnd.google-apps.presentation'
    MIMETYPE_MS_EXCEL = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    MIMETYPE_MS_POWERPOINT = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
    MIMETYPE_PDF = 'application/pdf'

    try:
        # --- Google Slides (Presentations) ---
        if mime_type == MIMETYPE_GOOGLE_PRESENTATION:
            print(f"    Exporting Google Slides '{file_name}' as PPTX for slide-level parsing...")
            # Export Google Slides as PPTX
            request_pptx = drive_service.files().export_media(
                fileId=file_id,
                mimeType=MIMETYPE_MS_POWERPOINT # Export as PPTX
            )
            fh_pptx = io.BytesIO()
            downloader_pptx = MediaIoBaseDownload(fh_pptx, request_pptx)
            done_pptx = False
            while not done_pptx:
                status, done_pptx = downloader_pptx.next_chunk()
            fh_pptx.seek(0)

            try:
                prs = Presentation(fh_pptx)
                slides_data = []
                for i, slide in enumerate(prs.slides):
                    slide_text_runs = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text_frame") and shape.text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if run.text:
                                        slide_text_runs.append(run.text)
                    slides_data.append({
                        "slide_number": i + 1,
                        "text": "\n".join(slide_text_runs),
                        "source_type": "Google Presentation (exported as PPTX)"
                    })
                if not slides_data and prs.slides: # Had slides but no text extracted
                     return f"Google Slides '{file_name}': Exported as PPTX, but no text content found in slides."
                elif not prs.slides: # Exported PPTX was empty
                    return f"Google Slides '{file_name}': Exported as PPTX, but it contained no slides."
                return slides_data                          
            except Exception as e_pptx_parse:
                print(f"    Warning: Failed to parse Google Slides '{file_name}' exported as PPTX: {e_pptx_parse}. Falling back to plain text export.")
                # Fallback: Export as plain text if PPTX export/parse fails
                request_text = drive_service.files().export_media(fileId=file_id, mimeType='text/plain')
                fh_text = io.BytesIO()
                downloader_text = MediaIoBaseDownload(fh_text, request_text)
                done_text = False
                while not done_text:
                    status, done_text = downloader_text.next_chunk()
                fh_text.seek(0)
                full_text = fh_text.read().decode('utf-8', errors='replace')
                if not full_text.strip():
                    return f"Google Slides '{file_name}': Fallback to plain text export, but no content found."
                return [{"slide_number": 1, "text": full_text, "source_type": "Google Presentation (all text fallback)"}]


        # --- Microsoft PowerPoint (.pptx) ---
        elif mime_type == MIMETYPE_MS_POWERPOINT: # Handles native PPTX files
            request = drive_service.files().get_media(fileId=file_id)
            fh = io.BytesIO()
            downloader = MediaIoBaseDownload(fh, request)
            done = False
            while not done: status, done = downloader.next_chunk()
            fh.seek(0)
            try:
                prs = Presentation(fh)
                slides_data = []
                for i, slide in enumerate(prs.slides):
                    slide_text_runs = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text_frame") and shape.text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if run.text: slide_text_runs.append(run.text)
                    slides_data.append({"slide_number": i + 1, "text": "\n".join(slide_text_runs), "source_type": "PPTX (native)"}) # Clarified source
                if not slides_data and prs.slides:
                    return f"PPTX File '{file_name}': Contained slides, but no text content found in them."
                elif not prs.slides:
                    return f"PPTX File '{file_name}': Contained no slides."
                return slides_data
            except Exception as e:
                return f"Could not parse native PPTX content for '{file_name}': {e}"
            

        # --- Google Sheets & Microsoft Excel ---
        elif mime_type == MIMETYPE_GOOGLE_SHEET or mime_type == MIMETYPE_MS_EXCEL:
            all_rows_data = []
            try:
                if mime_type == MIMETYPE_GOOGLE_SHEET:
                    spreadsheet = sheets_service.spreadsheets().get(spreadsheetId=file_id).execute()
                    for sheet_meta in spreadsheet.get('sheets', []):
                        sheet_title = sheet_meta['properties']['title']
                        # Read a significant number of rows, e.g., 1000, and all columns up to ZZ
                        range_str = f"'{sheet_title}'!A1:AZ5000" 
                        result = sheets_service.spreadsheets().values().get(
                            spreadsheetId=file_id, range=range_str
                        ).execute()
                        rows = result.get('values', [])
                        if rows:
                            header = rows[0]
                            for row_idx, row_values in enumerate(rows): # (row_idx is 0-based here)
                                all_rows_data.append({"sheet_name": sheet_title, 
                                                      "row_index": row_idx + 1, # 1-based for display
                                                      "header": header, 
                                                      "values": row_values,
                                                      "source_type": "Google Sheet"})
                elif mime_type == MIMETYPE_MS_EXCEL:
                    request = drive_service.files().get_media(fileId=file_id)
                    fh = io.BytesIO()
                    downloader = MediaIoBaseDownload(fh, request)
                    done = False
                    while not done: status, done = downloader.next_chunk()
                    fh.seek(0)
                    workbook = openpyxl.load_workbook(fh)
                    for sheet_title in workbook.sheetnames:
                        sheet = workbook[sheet_title]
                        header = [cell.value for cell in sheet[1]] # Assuming header is 1st row
                        for row_idx, row_obj in enumerate(sheet.iter_rows(min_row=1, max_col=50, max_row=1000, values_only=True)): # max_col to limit width
                            # row_idx is 0-based here from iter_rows(min_row=1)
                            all_rows_data.append({"sheet_name": sheet_title,
                                                  "row_index": row_idx + 1, # 1-based for display
                                                  "header": header, 
                                                  "values": list(row_obj), # Ensure it's a list
                                                  "source_type": "Excel Sheet"})
                if not all_rows_data: # If loops completed but no data (e.g. all sheets were empty)
                    return f"Spreadsheet file '{file_name}' ({mime_type}) processed, but no data rows found in any sheet."
                return all_rows_data
            except Exception as e:
                return f"Could not parse Spreadsheet content for '{file_name}' ({mime_type}): {e}"
        
        # --- Google Documents ---
        elif mime_type == MIMETYPE_GOOGLE_DOC:
            request = drive_service.files().export_media(fileId=file_id, mimeType='text/plain')
            fh = io.BytesIO()
            downloader = MediaIoBaseDownload(fh, request)
            done = False
            while not done: status, done = downloader.next_chunk()
            fh.seek(0)
            return fh.read().decode('utf-8', errors='replace') # Returns string

        # --- PDF Files ---
        elif mime_type == MIMETYPE_PDF:
            try:
                request = drive_service.files().get_media(fileId=file_id)
                fh = io.BytesIO()
                downloader = MediaIoBaseDownload(fh, request)
                done = False
                while not done:
                    status, done = downloader.next_chunk()
                fh.seek(0)
                pdf_text = ""
                with fitz.open(stream=fh.read(), filetype="pdf") as doc:
                    for page in doc:
                        pdf_text += page.get_text() + "\n"
                if not pdf_text.strip():
                    return f"PDF File: '{file_name}'. No extractable text found."
                return pdf_text
            except Exception as e:
                return f"Could not parse PDF content: {e}" 

        # --- Plain Text Files ---
        elif mime_type.startswith('text/'):
            request = drive_service.files().get_media(fileId=file_id)
            fh = io.BytesIO()
            downloader = MediaIoBaseDownload(fh, request)
            done = False
            while not done: status, done = downloader.next_chunk()
            fh.seek(0)
            return fh.read().decode('utf-8', errors='replace') # Returns string
            
        else:
            return f"File type ({mime_type}) for '{file_name}' not configured for structured data extraction. Its name and existence might be relevant." # Returns string

    except HttpError as error:
        return f"An HTTP error occurred reading GDrive file {file_name} (ID: {file_id}): {error}" # Returns string
    except Exception as e:
        return f"A general error occurred reading GDrive file {file_name} (ID: {file_id}): {e}" # Returns string


def summarize_file_content_with_gemini(gemini_llm_client, file_name, mime_type, file_content):
    """
    Uses Gemini LLM to summarize the content of a file for inclusion in the meeting brief.
    """
    if not gemini_llm_client:
        return "Error: Gemini model not available for summarization."

    prompt = (
        f"Summarize the following content from the file '{file_name}' (type: {mime_type}) in 5-10 concise bullet points, "
        "focusing on key facts, data, or insights that would be useful for a marketing/sales meeting. "
        "Do not include generic statements. If the content is not relevant, say 'No relevant content found.'\n\n"
        f"---\n{file_content}\n---"
    )
    try:
        response = gemini_llm_client.models.generate_content(model="gemini-2.5-flash",contents=prompt)
        if response and response.candidates and response.candidates[0].content.parts:
            summary = response.candidates[0].content.parts[0].text.strip()
            return summary
        else:
            return "No summary generated."
    except Exception as e:
        print(f"Error during Gemini summarization: {e}")
        return f"Error: Exception during Gemini summarization: {e}"





# --- Modify get_internal_nbh_data_for_brand ---
def get_internal_nbh_data_for_brand(drive_service, sheets_service, gemini_llm_client, 
                                    current_target_brand_name, target_brand_industry, current_meeting_data, 
                                    EXCLUDED_NBH_PSEUDO_NAMES_FOR_FOLLOWUP, AGENT_EMAIL, master_sheet_id):
    """
    Fetches internal data with STRICT Privacy Logic.
    1. Search Master Sheet for Brand History.
    2. Check for Attendee Overlap (At least one matching NBH person).
    3. IF MATCH: Extract Action Items & Discussions for the Brief.
    4. IF NO MATCH: Do NOT extract details (Privacy).
    """
    print(f"Fetching and processing internal NBH data for target brand '{current_target_brand_name}'...")
    
    # Containers
    history_context_str = ""
    general_context_parts = []
    
    # Flags
    is_overall_direct_follow_up = False
    has_other_past_interactions = False 
    condensed_past_meetings_for_alert = []

    # --- HELPER FUNCTION FOR CAMPAIGN EXTRACTION (RESTORED) ---
    def extract_matching_rows(file_data_obj, fname, brand_clean, kwords):
        list_2025 = []
        list_2024 = []
        
        if not isinstance(file_data_obj, list) or not file_data_obj:
            return list_2025, list_2024
        
        header_vals = file_data_obj[0].get("header", []) if isinstance(file_data_obj[0], dict) else []
        brand_col, ind_col, date_col = -1, -1, -1

        if header_vals:
            lower_h = [str(h).strip().lower() for h in header_vals]
            for idx, h in enumerate(lower_h):
                if "brand" in h: brand_col = idx
                if any(x in h for x in ["industry", "category", "vertical", "segment"]): ind_col = idx
                if any(x in h for x in ["year", "date", "timestamp", "month"]): date_col = idx

        if brand_col == -1:
            return list_2025, list_2024

        data_rows = file_data_obj[1:] if len(file_data_obj) > 1 else []
        
        for row_info in reversed(data_rows):
            vals = row_info.get('values', [])
            if not vals or not any(str(v).strip() for v in vals): 
                continue
            
            row_brand = str(vals[brand_col]).strip().lower() if len(vals) > brand_col else ""
            row_ind = str(vals[ind_col]).strip().lower() if ind_col != -1 and len(vals) > ind_col else ""
            row_date = str(vals[date_col]).strip() if date_col != -1 and len(vals) > date_col else ""

            is_match = False
            if brand_clean in row_brand or row_brand in brand_clean:
                is_match = True 
            elif kwords and any(word in row_ind for word in kwords):
                is_match = True 

            if is_match:
                row_items = [f"{header_vals[i]}: {str(vals[i]).strip()}" 
                            for i in range(len(header_vals)) 
                            if i < len(vals) and str(vals[i]).strip() and str(vals[i]).lower() != "n/a"]
                entry = " | ".join(row_items) + "\n"
                
                if "2025" in row_date:
                    list_2025.append(entry)
                elif "2024" in row_date:
                    list_2024.append(entry)
        
        return list_2025, list_2024

    # --- SETUP INDUSTRY KEYWORDS FOR CAMPAIGN SEARCH ---
    STRICT_INDUSTRY_MAP = {
        "FMCG": ["fmcg", "consumer", "goods", "staples", "packaged", "textiles"],
        "Automotive & Transportation": ["automotive", "car", "bike", "ev", "vehicle", "transport"],
        "Food & Beverage": ["food", "beverage", "f&b", "dairy", "snacks", "drinks", "restaurant", "hospitality"],
        "Jewellery": ["jewel", "gold", "diamond", "ornament"],
        "Apparel & Fashion": ["apparel", "fashion", "clothing", "wear", "shoes", "textiles"],
        "Finance & Fintech": ["finance", "fintech", "bank", "insurance", "loan", "trading"],
        "Beauty & Personal Care": ["beauty", "cosmetic", "skincare", "grooming", "wellness"],
        "Real Estate & Construction": ["real estate", "builder", "construction", "property", "interior design", "interior"],
        "Healthcare": ["healthcare", "hospital", "medical", "pharma", "wellness", "clinic"],
        "E-Commerce": ["e-commerce", "ecommerce", "online", "marketplace"],
        "Retail": ["retail", "supermarket", "mall", "store", "furniture"],
        "OTT": ["ott", "streaming", "entertainment", "video"],
        "Marketing, Advertising & Media": ["advertising", "marketing", "events", "entertainment", "media"],
        "Education & Training": ["education", "school", "college", "training", "admissions"],
        "Home Goods & Electronics": ["furniture", "interior", "home services", "electronics", "appliances"],
        "Hospitality & Travel": ["hospitality", "travel", "tourism", "hotel", "resort"],
        "Membership & Local Services": ["home services", "local", "membership", "community"]
    }
    keywords = STRICT_INDUSTRY_MAP.get(target_brand_industry, [])
    target_brand_clean = current_target_brand_name.lower().strip()

    # 1. Get Current Meeting Date
    current_meeting_date_obj = current_meeting_data.get('start_time_obj')
    if isinstance(current_meeting_date_obj, datetime.datetime):
        current_meeting_date_only = current_meeting_date_obj.date()
    else:
        current_meeting_date_only = datetime.date.today()

    # 2. Extract Current NBH Attendees (Names & Emails) for Matching
    current_nbh_tokens = set()
    for att in current_meeting_data.get('nbh_attendees', []):
        if att.get('email'):
            current_nbh_tokens.add(att['email'].lower().split('@')[0].strip()) 
        if att.get('name'):
             parts = att['name'].lower().split()
             for p in parts: 
                 if len(p) > 2: current_nbh_tokens.add(p)

    # 3. CHECK MASTER SHEET DIRECTLY
    print(f"    Checking Master Sheet for Brand History: '{current_target_brand_name}'...")
    try:
        # Fetch Headers
        header_req = sheets_service.spreadsheets().values().get(spreadsheetId=master_sheet_id, range="Meeting_data!A1:AZ1").execute()
        headers = header_req.get('values', [])[0]
        lower_headers = [str(h).strip().lower() for h in headers]

        # Map Columns
        try:
            col_brand = lower_headers.index("brand name")
            col_date = lower_headers.index("meeting date")
            col_discussion = lower_headers.index("key discussion points")
            col_actions = lower_headers.index("action items")
            col_nbh_attendees = lower_headers.index("nobroker attendees")
            # Optional
            col_pain = lower_headers.index("client pain points") if "client pain points" in lower_headers else -1
        except ValueError as e:
            print(f"    CRITICAL: Master Sheet missing required columns. {e}")
            history_context_str = "" # Fail safe
        else:
            # Fetch Data
            data_req = sheets_service.spreadsheets().values().get(spreadsheetId=master_sheet_id, range="Meeting_data!A2:AZ").execute()
            data_rows = data_req.get('values', [])

            found_meetings = []
            
            for row in data_rows:
                if len(row) <= col_brand: continue
                
                sheet_brand = str(row[col_brand]).strip()
                
                # A. Fuzzy Brand Match
                if current_target_brand_name.lower() in sheet_brand.lower() or sheet_brand.lower() in current_target_brand_name.lower():
                    
                    # B. Date Check (Past Only)
                    row_date_str = str(row[col_date]) if len(row) > col_date else ""
                    try:
                        clean_date = row_date_str.split(" ")[0]
                        row_date_obj = datetime.datetime.strptime(clean_date, "%Y-%m-%d").date()
                    except:
                        try: 
                            clean_date = row_date_str.split(" ")[0]
                            row_date_obj = datetime.datetime.strptime(clean_date, "%d/%m/%Y").date()
                        except: continue 
                    
                    if row_date_obj < current_meeting_date_only:
                        
                        # C. ATTENDEE MATCHING (Strict Privacy Check)
                        prev_nbh_raw = str(row[col_nbh_attendees]).lower() if len(row) > col_nbh_attendees else ""
                        
                        is_match = False
                        if prev_nbh_raw:
                            is_match = any(token in prev_nbh_raw for token in current_nbh_tokens)
                        
                        meeting_info = {
                            "date": row_date_obj,
                            "date_str": row_date_str,
                            "discussion": row[col_discussion] if len(row) > col_discussion else "N/A",
                            "actions": row[col_actions] if len(row) > col_actions else "None recorded",
                            "pain_points": row[col_pain] if col_pain != -1 and len(row) > col_pain else "N/A",
                            "nbh_team": prev_nbh_raw
                        }

                        if is_match:
                            found_meetings.append(meeting_info)
                        else:
                            has_other_past_interactions = True
                            condensed_past_meetings_for_alert.append({
                                "date": row_date_str,
                                "discussion_summary": "Different Team - Content Hidden",
                                "nbh_team": prev_nbh_raw
                            })

            # Sort Matched Meetings (Newest First)
            found_meetings.sort(key=lambda x: x['date'], reverse=True)

            if found_meetings:
                is_overall_direct_follow_up = True
                top_meeting = found_meetings[0]

                # --- CONSTRUCT THE "CALL PREP" INTELLIGENCE BLOCK ---
                history_context_str = (
                    f"## PREVIOUS MEETING INTELLIGENCE (MATCHED)\n"
                    f"**Last Meeting Date:** {top_meeting['date_str']}\n"
                    f"**Last NBH Team:** {top_meeting['nbh_team']}\n"
                    f"--------------------------------------------------\n"
                    f"**>>> LAST MEETING ACTION ITEMS (Review Required):**\n"
                    f"{top_meeting['actions']}\n\n"
                    f"**>>> KEY DISCUSSION SUMMARY:**\n"
                    f"{top_meeting['discussion']}\n\n"
                    f"**>>> CLIENT PAIN POINTS:**\n"
                    f"{top_meeting['pain_points']}\n"
                    f"--------------------------------------------------\n"
                )
                
                # Add depth from 2nd latest meeting
                if len(found_meetings) > 1:
                    history_context_str += "**Older Context:**\n"
                    for old in found_meetings[1:3]:
                        history_context_str += f"- {old['date_str']}: {old['discussion'][:200]}...\n"

    except Exception as e:
        print(f"    Error reading Master Sheet: {e}")

    # 4. PROCESS OTHER FILES (Pitch Decks, Campaigns)
    all_files_in_folder = list_files_in_gdrive_folder(drive_service, NBH_GDRIVE_FOLDER_ID)
    for item in all_files_in_folder:
        fname = item.get('name', '')
        fid = item['id']
        mtype = item.get('mimeType', '')

        if FILE_NAME_NBH_PREVIOUS_MEETINGS_GSHEET.lower() in fname.lower(): continue

        # Campaigns/Case Studies using the restored helper function
        if FILE_NAME_PHYSICAL_CAMPAIGNS_GSHEET.lower() in fname.lower():
            content = get_structured_gdrive_file_data(drive_service, sheets_service, fid, fname, mtype)
            p_25, p_24 = extract_matching_rows(content, fname, target_brand_clean, keywords)
            if p_25 or p_24:
                general_context_parts.append(f"## PHYSICAL CAMPAIGNS ({fname}):\n" + "\n".join(p_25[:5] + p_24[:3]))

        elif FILE_NAME_DIGITAL_CAMPAIGNS_GSHEET.lower() in fname.lower():
            content = get_structured_gdrive_file_data(drive_service, sheets_service, fid, fname, mtype)
            d_25, d_24 = extract_matching_rows(content, fname, target_brand_clean, keywords)
            if d_25 or d_24:
                general_context_parts.append(f"## DIGITAL CAMPAIGNS ({fname}):\n" + "\n".join(d_25[:5] + d_24[:3]))

        elif FILE_NAME_LATEST_CASE_STUDIES_GSHEET.lower() in fname.lower():
             pass 
        
        elif 'pdf' in mtype.lower() and (FILE_NAME_PITCH_DECK_PDF.lower() in fname.lower() or FILE_NAME_CASE_STUDIES_PDF.lower() in fname.lower()):
             content = get_structured_gdrive_file_data(drive_service, sheets_service, fid, fname, mtype)
             if isinstance(content, str) and len(content) > 100:
                 general_context_parts.append(f"## Reference Document: {fname}\n{content[:5000]}...\n")

    # 5. FINAL ASSEMBLY
    final_llm_string = f"{history_context_str}\n\n" + "\n".join(general_context_parts)
    
    return {
        "llm_summary_string": final_llm_string,
        "is_overall_direct_follow_up": is_overall_direct_follow_up,
        "has_other_past_interactions": has_other_past_interactions,
        "condensed_past_meetings_for_alert": condensed_past_meetings_for_alert
    }

# --- Calendar Processing ---
def get_upcoming_meetings(calendar_service, calendar_id='primary', time_delta_hours=96): # Process meetings in next 3 days
    now_utc = datetime.datetime.utcnow()
    time_min_str = (now_utc - datetime.timedelta(hours=48)).isoformat() + 'Z'
    time_max_str = (now_utc + datetime.timedelta(hours=time_delta_hours)).isoformat() + 'Z'
    
    print(f'Getting events between {time_min_str} and {time_max_str}')
    try:
        events_result = calendar_service.events().list(
            calendarId=calendar_id, timeMin=time_min_str, timeMax=time_max_str,
            singleEvents=True, orderBy='startTime',
            # q='brand.vmeet@nobroker.in' # This might filter too early if brandvmeet is added as resource
        ).execute()
        events = events_result.get('items', [])
        
        # ✅ FILTER OUT TASK EVENTS (Skip events starting with task markers)
        TASK_PREFIXES = ['✅ TASK', '☑ TASK', 'TASK:', '[TASK]']
        filtered_events = []
        
        for event in events:
            title = event.get('summary', '').strip()
            is_task = any(title.startswith(prefix) for prefix in TASK_PREFIXES)
            
            if not is_task:
                filtered_events.append(event)
            else:
                print(f"  ⏭️  Skipping task event: '{title}'")
        
        print(f"  ✅ Filtered {len(events) - len(filtered_events)} task events, {len(filtered_events)} meetings remaining")
        return filtered_events
        
    except HttpError as error:
        print(f'An error occurred fetching events: {error}')
        return []

def load_processed_event_ids():
    # ... (same as before) ...
    if not os.path.exists(PROCESSED_EVENTS_FILE): return set()
    with open(PROCESSED_EVENTS_FILE, 'r') as f: return set(line.strip() for line in f)

def save_processed_event_id(event_id):
    # ... (same as before) ...
    with open(PROCESSED_EVENTS_FILE, 'a') as f: f.write(event_id + '\n')

EVENT_TAG_PROCESSED = "[NBH_BRIEF_AGENT_PROCESSED_V1]"

def is_event_already_tagged(event_description):
    return EVENT_TAG_PROCESSED in (event_description or "")

def tag_event_as_processed(calendar_service, event_id, calendar_id='primary'):
    if not calendar_service:
        print("  Calendar service not available to tag event.")
        return
    try:
        event = calendar_service.events().get(calendarId=calendar_id, eventId=event_id).execute()
        description = event.get('description', '')
        if not is_event_already_tagged(description):
            new_description = f"{description}\n\n{EVENT_TAG_PROCESSED}"
            updated_event_body = {'description': new_description}
            calendar_service.events().patch(
                calendarId=calendar_id, eventId=event_id, body=updated_event_body
            ).execute()
            print(f"  Tagged event {event_id} as processed in calendar.")
    except HttpError as error:
        print(f"  An error occurred tagging event {event_id}: {error}")

EVENT_REMINDER_SET_TAG = "[NBH_1HR_REMINDER_SET]"

def set_one_hour_email_reminder(calendar_service, event_id, calendar_id='primary'):
    if not calendar_service:
        print("  Calendar service not available to set reminder.")
        return False
    try:
        event = calendar_service.events().get(calendarId=calendar_id, eventId=event_id).execute()
        
        # Check if our specific reminder tag is already there
        description = event.get('description', '')
        if EVENT_REMINDER_SET_TAG in description:
            print(f"  1-hour email reminder already marked as set for event {event_id}.")
            return True

        reminders = event.get('reminders', {})
        overrides = reminders.get('overrides', [])
        
        # Check if a 60-minute email reminder already exists
        has_one_hour_email_reminder = any(
            r.get('method') == 'email' and r.get('minutes') == 60 for r in overrides
        )

        if has_one_hour_email_reminder:
            print(f"  Event {event_id} already has a 60-minute email reminder.")
        else:
            print(f"  Adding 60-minute email reminder to event {event_id}.")
            overrides.append({'method': 'email', 'minutes': 60})
            body_update = {'reminders': {'useDefault': False, 'overrides': overrides}}
            calendar_service.events().patch(
                calendarId=calendar_id, eventId=event_id, body=body_update
            ).execute()
            print(f"  Successfully added 60-minute email reminder for event {event_id}.")

        # Add our custom tag to the description to indicate we've processed this for reminders
        if EVENT_REMINDER_SET_TAG not in description:
            new_description = f"{description}\n{EVENT_REMINDER_SET_TAG}".strip()
            body_update_desc = {'description': new_description}
            calendar_service.events().patch(
                calendarId=calendar_id, eventId=event_id, body=body_update_desc
            ).execute()
            print(f"  Tagged event {event_id} with {EVENT_REMINDER_SET_TAG}.")
        
        return True

    except HttpError as error:
        print(f"  An error occurred setting reminder for event {event_id}: {error}")
        return False



# --- Meeting Info Extraction ---

NBH_SERVICE_ACCOUNTS_TO_EXCLUDE = { # Emails to exclude from the displayed NBH attendee list
    AGENT_EMAIL.lower(),
    "pia.brand@nobroker.in",
    "pia@nobroker.in",
    "nbh.meeting@gmail.com"
}


def extract_meeting_info(event, agent_email_global, nbh_service_accounts_to_exclude_global):
    """
    Extracts basic, non-inferential information from a calendar event.
    The brand name itself is NOT processed here; it is extracted by the LLM later.
    """
    event_id = event['id']
    summary = event.get('summary', 'No Title')
    start_str = event['start'].get('dateTime', event['start'].get('date'))
    start_time_obj = datetime.datetime.fromisoformat(start_str.replace('Z', '+00:00'))
    location = event.get('location', 'N/A')
    description = event.get('description', '')
    attendees = event.get('attendees', [])

    nbh_attendees = []
    brand_attendees_info = []

    is_agent_invited = any(attendee.get('email', '').lower() == agent_email_global.lower() for attendee in attendees)
    if not is_agent_invited:
        print(f"  Skipping event '{summary}': {agent_email_global} is not an attendee.")
        return None

    for attendee in attendees:
        email = attendee.get('email', '').lower()
        name = attendee.get('displayName', email.split('@')[0] if '@' in email else email)
        
        # Skip if this is an excluded service account (check BEFORE categorization)
        if email in nbh_service_accounts_to_exclude_global:
            continue
        
        # Categorize as NBH or Brand attendee
        if '@nobroker.in' in email:
            nbh_attendees.append({'email': email, 'name': name})
        elif email:
            brand_attendees_info.append({'name': name, 'email': email})
    # Removing this condition so that physical meetings do not get skipped
    # if not brand_attendees_info:
    #     print(f"  Skipping event '{summary}': No external attendees.")
    #     return "NO_EXTERNAL_ATTENDEES"

    return {
        'id': event_id,
        'title': summary, # Return the raw title
        'start_time_obj': start_time_obj,
        'start_time_str': start_time_obj.strftime("%Y-%m-%d %I:%M %p %Z (%A)"),
        'location': location,
        'description': description,
        'nbh_attendees': nbh_attendees,
        'brand_attendees_info': brand_attendees_info,
        'is_event_description_present_for_tagging': bool(event.get('description'))
    }


# --- Gemini LLM Integration ---
def configure_gemini():
    """
    Configures and returns a Gemini LLM model instance using the provided API key.
    
    Returns:
        A configured Gemini GenerativeModel object if successful, or None if configuration fails or the API key is missing.
    """
    if not GEMINI_API_KEY:
        print("GEMINI_API_KEY environment variable not set. LLM will not function.")
        return None
    try:
        client = genai.Client(api_key=GEMINI_API_KEY)
        # Using a specific model version. 1.5 Flash is faster and cheaper for many tasks.
        # For higher quality, consider 'gemini-1.5-pro-latest'.
        # model = genai.GenerativeModel('gemini-2.5-flash')
        print(f"Gemini model configured successfully.")
        return client
    except Exception as e:
        print(f"Error configuring Gemini API: {e}")
        return None  


def generate_brief_with_gemini(gemini_llm_client, YOUR_DETAILED_PROMPT_TEMPLATE_GEMINI, meeting_data, internal_data_summary_str):
    if not gemini_llm_client:
        return "Error: Gemini model not available."

    nbh_attendee_names_str = ", ".join([att['name'] for att in meeting_data['nbh_attendees']])
    brand_attendee_names_only_str = ", ".join([att['name'] for att in meeting_data['brand_attendees_info']])
    
    # ========== NEW CODE: Format brand attendees with LinkedIn URLs ==========
    brand_attendees_with_linkedin_str = ""
    for att in meeting_data['brand_attendees_info']:
        linkedin_display = att.get('linkedin_url', '(LinkedIn Not Verified)')
        
        # If we have a real URL, make it a markdown link
        if linkedin_display and linkedin_display != '(LinkedIn Not Verified)':
            linkedin_display = f"[LinkedIn Profile]({linkedin_display})"
        
        # Add this attendee's info to the string
        brand_attendees_with_linkedin_str += f"- **{att['name']}** ({att['email']}) - {linkedin_display}\n"
    
    # Keep the old format too for backward compatibility
    brand_attendees_info_str = "; ".join([f"{att['name']} ({att['email']})" for att in meeting_data['brand_attendees_info']])
    # ========== END NEW CODE ==========

   # ========== NEW CODE: Format potential key contacts ==========
    potential_contacts_str = ""
    key_contacts_list = meeting_data.get('potential_key_contacts', [])

    if key_contacts_list:
        potential_contacts_str = "**Found Key Contacts:**\n\n"
        for contact in key_contacts_list:
            linkedin_display = contact.get('linkedin_url', '(LinkedIn Not Verified)')
            if linkedin_display and linkedin_display != '(LinkedIn Not Verified)':
                linkedin_display = f"[LinkedIn Profile]({linkedin_display})"
            
            potential_contacts_str += f"- **{contact['name']}** - {contact['title']} - {linkedin_display}\n"
            potential_contacts_str += f"  - Why They Matter: {contact['reasoning']}\n\n"
    else:
        potential_contacts_str = "**No additional key contacts found through search.**\n\n"
    # ========== END NEW CODE ==========


    prompt_filled = YOUR_DETAILED_PROMPT_TEMPLATE_GEMINI.format(
    MEETING_DATETIME=meeting_data['start_time_str'],
    MEETING_LOCATION=meeting_data['location'],
    BRAND_NAME=meeting_data['brand_name'],
    BRAND_ATTENDEES_NAMES=brand_attendee_names_only_str,
    NBH_ATTENDEES_NAMES=nbh_attendee_names_str,
    BRAND_NAME_FOR_BODY=meeting_data['brand_name'],
    MEETING_TITLE=meeting_data.get('title', 'N/A'),
    BRAND_ATTENDEES_FULL_DETAILS=brand_attendees_info_str,
    BRAND_ATTENDEES_WITH_LINKEDIN=brand_attendees_with_linkedin_str,
    POTENTIAL_KEY_CONTACTS=potential_contacts_str,
    INTERNAL_NBH_DATA_SUMMARY=internal_data_summary_str
)
    
    grounding_tool = types.Tool(
        google_search=types.GoogleSearch()
    )

    # Configure generation settings
    config = types.GenerateContentConfig(
    # sampling parameters (formerly generation_config dict)
    temperature=0.0,
    top_p=0.95,
    top_k=40,

    # safety filters (formerly safety_settings list of dicts)
    safety_settings=[
        types.SafetySetting(
            category=types.HarmCategory.HARM_CATEGORY_HARASSMENT,
            threshold=types.HarmBlockThreshold.BLOCK_MEDIUM_AND_ABOVE,
        ),
        types.SafetySetting(
            category=types.HarmCategory.HARM_CATEGORY_HATE_SPEECH,
            threshold=types.HarmBlockThreshold.BLOCK_MEDIUM_AND_ABOVE,
        ),
        types.SafetySetting(
            category=types.HarmCategory.HARM_CATEGORY_SEXUALLY_EXPLICIT,
            threshold=types.HarmBlockThreshold.BLOCK_MEDIUM_AND_ABOVE,
        ),
        types.SafetySetting(
            category=types.HarmCategory.HARM_CATEGORY_DANGEROUS_CONTENT,
            threshold=types.HarmBlockThreshold.BLOCK_MEDIUM_AND_ABOVE,
        ),
    ],

    # tools (formerly a separate config.tools or tools argument)
    tools=[grounding_tool],
)

    print(f"  Sending request to Gemini for brand: {meeting_data['brand_name']}...")
    try:
        # Use Google Search grounding by adding tools=[{"tool": "google_search"}]
        response = gemini_llm_client.models.generate_content(
            model="gemini-2.5-flash",  # Use the latest Gemini model
            contents=prompt_filled,
            config=config,
            # tools=[{"tool": "google_search"}]  # Enable Google Search grounding >> not working on API calls rights now
        )
        # print(f"Gemini raw response: {response}") # For debugging
        if response.prompt_feedback and response.prompt_feedback.block_reason:
            return f"Error: Prompt blocked by Gemini. Reason: {response.prompt_feedback.block_reason_message or response.prompt_feedback.block_reason}"

        if not response.candidates:
             return f"Error: Gemini returned no candidates. Feedback: {response.prompt_feedback}"

        # Ensure text is extracted correctly
        brief_content = ""
        for part in response.candidates[0].content.parts:
            brief_content += part.text
        
        if not brief_content.strip():
            return "Error: Gemini returned an empty brief."
            
        return brief_content
    except Exception as e:
        print(f"  Error during Gemini API call: {e}")
        # Check for specific API errors if needed, e.g. google.api_core.exceptions.ResourceExhausted
        if "429" in str(e) or "ResourceExhausted" in str(e): # Quota issue
             return "Error: Gemini API quota likely exceeded. Please check your Google Cloud/AI Studio quotas."
        return f"Error: Exception during Gemini call: {e}"

# --- Email Sending ---
def create_email_message(sender, to_emails_list, subject, message_text_html):
    # ... (Using MIMEText with 'html' for better formatting) ...
    message = MIMEText(message_text_html, 'html', 'utf-8') # Send as HTML
    message['to'] = ", ".join(to_emails_list)
    message['from'] = sender
    message['subject'] = subject
    raw_message = base64.urlsafe_b64encode(message.as_bytes())
    return {'raw': raw_message.decode()}

def send_gmail_message(gmail_service, user_id, message_body):
    # ... (same as before) ...
    """
    Sends an email message using the Gmail API.
    
    Parameters:
        user_id (str): The user's email address or the special value 'me' to indicate the authenticated user.
        message_body (dict): The message payload, typically created by create_email_message().
    
    Returns:
        dict or None: The API response containing the sent message details, or None if sending fails.
    """
    if not gmail_service:
        print("  Gmail service not available. Cannot send email.")
        return None
    try:
        message = (gmail_service.users().messages().send(userId=user_id, body=message_body).execute())
        print(f'  Message Id: {message["id"]} sent.')
        return message
    except HttpError as error:
        print(f'  An error occurred sending email: {error}')
        return None

def send_brief_email(gmail_service, meeting_data, brief_content):
    """
    Sends a pre-meeting brief email to NBH internal attendees for a given meeting.
    
    The function filters out excluded emails from the NBH attendee list, converts the brief content from markdown to HTML, formats the email body, and sends the email using the Gmail API. If no eligible recipients are found, the function exits without sending an email.
    """
    EXCLUDED_EMAILS = {AGENT_EMAIL.lower(), "pia.brand@nobroker.in","pia.hood@nobroker.in"} # Define a set of excluded emails

    nbh_recipient_emails = []
    attendees_list = meeting_data.get('nbh_attendees', []) 
    if isinstance(attendees_list, list): # Extra safety check
        for att in attendees_list:
            # Ensure 'att' is a dictionary and 'email' key exists
            if isinstance(att, dict) and 'email' in att:
                attendee_email = att.get('email')
                if attendee_email and isinstance(attendee_email, str) and attendee_email.lower() not in EXCLUDED_EMAILS:
                    nbh_recipient_emails.append(attendee_email)
    
    # For testing, override recipients:
    # nbh_recipient_emails = [ADMIN_EMAIL_FOR_NOTIFICATIONS]
    # print(f"DEBUG: Intended brief recipients: {nbh_recipient_emails}")

    if not nbh_recipient_emails:
        print(f"  No NBH recipients (other than brandvmeet) for '{meeting_data['title']}'. Brief not emailed.")
        return

    email_subject = f"Pre-Meeting Brief: {meeting_data['title']} with {meeting_data['brand_name']}"
    
    # Convert markdown-like brief (from LLM) to basic HTML for email
    html_brief_content = markdown.markdown(brief_content)
    # html_brief_content = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', html_brief_content) # Basic bold
    # Add more markdown to HTML conversions if needed (e.g., for lists, headers)  
    # For headers like # Detail Report, ## Sub-header
    # html_brief_content = html_brief_content.replace("# ", "<h3>").replace("\n", "</h3><br>") # Simple h3 for main headers
    
    email_body_html = f"""
    <html><body>
    <p>Hello Team,</p>
    <p>Please find the pre-meeting brief for your upcoming meeting:</p>
    <hr>
    {html_brief_content}
    <hr>
    <p>Best regards,<br>NBH Meeting Prep Agent</p>
    </body></html>
    """
    email_message = create_email_message(
        sender=AGENT_EMAIL,
        to_emails_list=nbh_recipient_emails,
        subject=email_subject,
        message_text_html=email_body_html
    )
    print(f"  FINAL CHECK: Sending brief for '{meeting_data['title']}' TO: {nbh_recipient_emails} FROM: {AGENT_EMAIL}")
    send_gmail_message(gmail_service, 'me', email_message)

def send_notification_email(gmail_service, subject, body_html, recipient=ADMIN_EMAIL_FOR_NOTIFICATIONS):
    if not recipient:
        print("  Admin notification email not set. Skipping notification.")
        return
    
    # Always send a copy to brandvmeet for record-keeping
    recipients = list(set([recipient, AGENT_EMAIL]))

    email_message = create_email_message(
        sender=AGENT_EMAIL,
        to_emails_list=recipients,
        subject=subject,
        message_text_html=body_html
    )
    send_gmail_message(gmail_service, 'me', email_message)

# Here are some useful functions to update meeting data in the master sheet

def read_data_from_sheets(sheet_id, sheets_service, range):

    try:
        result = (
                sheets_service.spreadsheets()
                .values()
                .get(spreadsheetId=sheet_id, range=range)
                .execute()
            )
        sheet_data = result.get("values", [])
        print(f"{len(sheet_data)} rows retrieved")
        return sheet_data
    except HttpError as error:
        print(f"An error occurred: {error}")

# Then check which events have not been updated in the google sheets
def events_to_update(meeting_ids, events):
    events_to_update = []
    for event in events:
        arr = [event["id"]]
        if arr not in meeting_ids:
            events_to_update.append(event)
    if not events_to_update:
        print("No new meetings to update")
        return None
    else:
        return events_to_update
    

def update_events_in_sheets(sheet_id, events_to_update, sheets_service, excluded_emails, designations):

    meeting_ids = read_data_from_sheets(sheet_id, sheets_service, "Meeting_data!A2:A")
    master_sheet_columns = read_data_from_sheets(sheet_id, sheets_service, "Meeting_data!A1:BZ1")[0]  # Get the header row
    audit_sheet_columns = read_data_from_sheets(sheet_id, sheets_service, "Audit_and_Training!A1:BZ1")[0]  # Get the header row
    owner_column_index_master = column_index[f"{master_sheet_columns.index('Owner') + 1}"]  # Convert to 1-based index
    owner_update_column_index_master = column_index[f"{master_sheet_columns.index('Owner sheet to be updated') + 1}"]  # Convert to 1-based index
    owner_column_index_audit = column_index[f"{audit_sheet_columns.index('Owner') + 1}"]  # Convert to 1-based index
    owner_update_column_index_audit = column_index[f"{audit_sheet_columns.index('Owner sheet to be updated') + 1}"]
    last_index = len(meeting_ids) + 1  # Start appending from the next row
    sheet_index = last_index + 1  # Sheet index starts from 1, so
    def to_rowdata(py_row):
        """Convert a list like ['Alice', 42] to RowData JSON."""
        def cell(v):
            t = "stringValue" if isinstance(v, str) else "numberValue"
            return {"userEnteredValue": {t: v}}
        return {"values": [cell(v) for v in py_row]}

    for i, event in enumerate(events_to_update):
        id = event["id"]
        # Safely get the event title. If it doesn't exist, use "Untitled Meeting".
        title = event.get("summary", "Untitled Meeting")
        date = event["start"].get("date", event["start"].get("dateTime"))
        if 'T' in date:
            date = datetime.datetime.fromisoformat(date).date().isoformat()
        attendees = event.get("attendees")
        if attendees:
            emails = [attendee["email"] for attendee in attendees]
            nobroker_attendee = []
            client_attendee = []
            
            for email in emails:
                # Skip excluded emails first (before categorization)
                if email.lower() in excluded_emails:
                    continue
                
                # Now categorize the remaining emails
                if "nobroker" in email:
                    nobroker_attendee.append(email)
                else:
                    client_attendee.append(email)
            row = [id, title, date, f"{nobroker_attendee}", f"{client_attendee}"]
            values = [to_rowdata(row)]
            try:
                requests = [
                    {
                        "appendCells": {
                            "sheetId": 0,  # Assuming Meeting_data is the first sheet
                            "rows": values,
                            "fields": "userEnteredValue"
                            }
                            },
                            {
                                "appendCells": {
                                    "sheetId": 1404820187,
                                    "rows": values,
                                    "fields": "userEnteredValue"
                                    }
                                    }
                                    ]
                body = {
                    "requests": requests
                }
                result = (
                    sheets_service.spreadsheets()
                    .batchUpdate(spreadsheetId=sheet_id, body=body)
                    .execute()
                    )
                print(f"Appended row: {title}")
                owner = None
                hierarchy_chain = []
                for email in nobroker_attendee:
                    owner, hierarchy_chain = get_sheet_owner_from_email(email)
                    if owner:
                        print(f"  Owner for {email} is {owner}. Hierarchy chain: {hierarchy_chain}")
                        break
                if not owner:
                    print(f"  No owner found for {title}")
                
                # Set the owner in the second last column and "True" value in the last column
                if owner:
                    values = [[owner, "TRUE"]]
                    update_body = {
                        "valueInputOption": 'USER_ENTERED',  # Use USER_ENTERED to allow date formatting
                        "data": [
                            {"range": f"Meeting_data!{owner_column_index_master}{sheet_index + i}:{owner_update_column_index_master}{sheet_index + i}", "values": values},
                            {"range": f"Audit_and_Training!{owner_column_index_audit}{sheet_index + i}:{owner_update_column_index_audit}{sheet_index + i}", "values": values},
                            ],
                        }
                    try:
                        resp = (
                            sheets_service.spreadsheets()
                            .values()
                            .batchUpdate(spreadsheetId=sheet_id, body=update_body)
                            .execute()
                            )
                        print(f"Updated owner and processed status for {title}")
                    except HttpError as error:
                        print(f"An error occurred while updating owner for {title}: {error}")
                # Updating the main participant and designation in the sheet; Main participant is generally the last person in the hierarchy chain. So if a BM is present, it will be the last person in the hierarchy chain. Else if RM is present, it will be the last person in the hierarchy chain. Else if CH is present, it will be the last person in the hierarchy chain.
                main_participant = []
                dg = []
                main_participant_column_index = column_index[f"{master_sheet_columns.index('Main participant') + 1}"]  # Convert to 1-based index
                meeting_done_status_column_index = column_index[f"{master_sheet_columns.index('Meeting Done') + 1}"]  # Convert to 1-based index
                for p in nobroker_attendee:
                    d = designations.get(p, None)
                    if d:
                        if d .lower() == 'bm':
                            main_participant.append(p)
                            dg.append(d)
                if not main_participant:
                    for p in nobroker_attendee:
                        d = designations.get(p, None)
                        if d:
                            if d.lower() == 'rm':
                                main_participant.append(p)
                                dg.append(d)
                if not main_participant:
                    for p in nobroker_attendee:
                        d = designations.get(p, None)
                        if d:
                            if d.lower() == 'ch':
                                main_participant.append(p)
                                dg.append(d)
                data = [[f"{main_participant}", f"{dg}", "Not Conducted"]]
                rng = f"Meeting_data!{main_participant_column_index}{sheet_index + i}:{meeting_done_status_column_index}{sheet_index + i}"
                values = data
                body = {
                    'values': values
                }
                try:
                    result = sheets_service.spreadsheets().values().update(
                        spreadsheetId=sheet_id,
                        range=rng,
                        valueInputOption='USER_ENTERED',
                        body=body
                    ).execute()
                    print(f"Updated main participant and designation for {title}")
                except HttpError as error:
                    print(f"An error occurred while updating main participant for {title}: {error}")
                
                # Sleep after every 50 updates to avoid rate limiting                    
                if (i+1)%50 == 0:
                    print("Sleep initiated")
                    time.sleep(50)
            except HttpError as error:
                print(f"An error occurred while updating {title}: {error}")
                continue

# Function to create a google document for the brief
def create_google_doc_in_folder(drive_service, folder_id, doc_name):
    file_metadata = {
        'name': doc_name,
        'mimeType': 'application/vnd.google-apps.document',
        'parents': [folder_id]
    }
    created = drive_service.files().create(
        body=file_metadata,
        fields='id, name, parents'
    ).execute()
    print(f"Created Google Doc: {created['name']} (ID: {created['id']})")
    return created['id']

# Function to write content to a Google Doc
def write_into_doc(docs_service, doc_id, text):
    requests = [
        {
            'insertText': {
                'location': { 'index': 1 },
                'text': text
            }
        }
    ]
    
    try:
        docs_service.documents().batchUpdate(
            documentId=doc_id,
            body={'requests': requests}
        ).execute()
    except:
        print("An error occured while writing into google doc")


def get_sheet_owner_from_email(email):
    hcy = []
    if email in sheet_masters:
        owner = email
        hcy.append(owner)
        return owner, hcy
    if email in hierarchy:
        manager = hierarchy[email]
        owner, hcy = get_sheet_owner_from_email(manager)
        hcy.append(email)
        return owner, hcy
    if email not in hierarchy:
        return None, []

# --- Main Execution Logic ---
def main():
    """
    Main orchestration function for automated pre-meeting brief generation and notification.
    
    This function coordinates the end-to-end workflow for preparing and emailing pre-meeting briefs for upcoming client meetings managed by the agent account. It initializes required Google Workspace services and the Gemini LLM, fetches upcoming calendar events, and processes each event as follows:
    
    - Skips events already processed or tagged.
    - Extracts meeting details and identifies the brand and industry using the LLM.
    - Retrieves and summarizes relevant internal NBH data for the brand.
    - Determines if the meeting is a direct follow-up or involves separate historical threads, and sends leadership alert emails as needed.
    - Generates a detailed pre-meeting brief using the LLM and internal data.
    - Emails the brief to NBH attendees, tags the event as processed, sets a 1-hour reminder, and records the event as processed.
    
    Handles error conditions gracefully, including missing services, ambiguous brand extraction, and LLM failures, with appropriate notifications and fallback logic.
    """
    print(f"Script started at {datetime.datetime.now()}")
    print(f"Using NBH GDrive Folder ID: {NBH_GDRIVE_FOLDER_ID}")
    
    # Load environment variables
    master_sheet_id = "1xtB1KUAXJ6IKMQab0Sb0NJfQppCKLkUERZ4PMZlNfOw"
    calendar_token = os.getenv("CALENDAR_TOKEN")
    gmail_token = os.getenv("GMAIL_TOKEN")
    drive_token = os.getenv("DRIVE_TOKEN")
    sheets_token = os.getenv("SHEET_TOKEN")
    docs_token = os.getenv("DOCS_TOKEN")  # Token for Google Docs API

    # Initialize Google Services
    # Use a combined token file strategy or separate ones. Separate is fine.
    calendar_service = get_google_service('calendar', 'v3', SCOPES, calendar_token)
    gmail_service = get_google_service('gmail', 'v1', SCOPES, gmail_token)
    drive_service = get_google_service('drive', 'v3', SCOPES, drive_token)
    sheets_service = get_google_service('sheets', 'v4', SCOPES, sheets_token)
    docs_service = get_google_service('docs', 'v1', SCOPES, docs_token)  # Docs service for creating briefs
    gemini_llm_client = configure_gemini()

    # Fetching employees data
    hcy_sheet_id = '1HxJt35QHF8BB_I8HusPQuiCS5_IpkEm5zoOSu1kwkNw'
    hcy_data = read_data_from_sheets(hcy_sheet_id, sheets_service, "Sheet4!A:F")
    df_hcy = pd.DataFrame(hcy_data[1:], columns=hcy_data[0])

    # Constructing designations dictionary
    # This will map employee emails to their designations
    designations = {}

    for i, row in df_hcy.iterrows():
        employee = row["Official Email ID"]
        dg = row["Designation New"]
        designations[employee] = dg

    # Fetching column headers for master sheet and audit sheet
    master_sheet_columns = read_data_from_sheets(master_sheet_id, sheets_service,  "Meeting_data!A1:BZ1")[0]  # Get the header row
    audit_sheet_columns = read_data_from_sheets(master_sheet_id, sheets_service, "Audit_and_Training!A1:BZ1")[0]
    # Create a mapping of column names to their 1-based index
    global column_index_master
    global column_index_audit
    column_index_master = {name: column_index[f"{i+1}"] for i, name in enumerate(master_sheet_columns)}
    column_index_audit = {name: column_index[f"{i+1}"] for i, name in enumerate(audit_sheet_columns)}

    prompts_sheet_id = "1_dKfSF_WkANgSNvFbMTR43By_sK74XKWUr9fTzire5s"
    pre_meeting_brief = "Pre_meeting_brief"
    rng = f"{pre_meeting_brief}!A2:A2"
    pre_meeting_brief_prompt = read_data_from_sheets(prompts_sheet_id, sheets_service, rng)
    YOUR_DETAILED_PROMPT_TEMPLATE_GEMINI = pre_meeting_brief_prompt[0][0]


    if not calendar_service: # Critical service
        print("Exiting: Calendar service failed to initialize.")
        return

    upcoming_events = get_upcoming_meetings(calendar_service)
    if not upcoming_events:
        print('No upcoming events found for agent email that need processing.')
        return

    # Updating events in the master sheet

    meeting_ids = read_data_from_sheets(master_sheet_id, sheets_service, "Meeting_data!A2:A")

    events_to_update_list = events_to_update(meeting_ids, upcoming_events)

    if not events_to_update_list:
        print("No new meetings to update in master sheet.")
    else:
        print(f"{len(events_to_update_list)} new meetings found")
        update_events_in_sheets(master_sheet_id, events_to_update_list, sheets_service, NBH_SERVICE_ACCOUNTS_TO_EXCLUDE, designations)
    

    updated_meeting_ids = read_data_from_sheets(master_sheet_id, sheets_service, "Meeting_data!A2:A")

    processed_ids_local_file = load_processed_event_ids()

    for event_payload in upcoming_events:
        event_id = event_payload['id']
        event_summary = event_payload.get('summary', 'No Title')
        event_description_for_tag_check = event_payload.get('description')

        print(f"\nProcessing event: '{event_summary}' (ID: {event_id})")

        # Step 1: Check if the event has already been processed
        if is_event_already_tagged(event_description_for_tag_check):
            print(f"  Skipping event '{event_summary}': Already tagged as processed.")
            continue
        
        if event_id in processed_ids_local_file:
            print(f"  Skipping event '{event_summary}': Found in local processed file.")
            continue

        # Step 2: Extract basic meeting info (attendees, raw title, etc.)
        meeting_data_result = extract_meeting_info(event_payload, AGENT_EMAIL, NBH_SERVICE_ACCOUNTS_TO_EXCLUDE)

        # Step 3: Handle the possible "skip" results from the extraction
        if meeting_data_result is None: # Case where agent is not an attendee
            print(f"  Skipping event '{event_summary}': Agent is not an attendee.")
            save_processed_event_id(event_id)
            tag_event_as_processed(calendar_service, event_id)
            continue
        # Skipping this condition so that physical meetings can also be processed
        # if meeting_data_result == "NO_EXTERNAL_ATTENDEES":
        #     print(f"  Event '{event_summary}': No external attendees. No brief needed.")
        #     save_processed_event_id(event_id)
        #     tag_event_as_processed(calendar_service, event_id)
        #     continue

        # Step 4: If we are here, extraction was successful. Assign the result to meeting_data.
        # This is the key fix: assign the dictionary before trying to use it.
        meeting_data = meeting_data_result

        # Step 5: Use the LLM to get the brand name and industry from the raw title
        print(f"  Using LLM to extract brand details from title: '{meeting_data['title']}'")
        
        # --- NEW SAFETY BLOCK FOR API QUOTA CRASHES ---
        try:
            brand_details = get_brand_details_from_title_with_llm(gemini_llm_client, meeting_data['title'])
            # Add a pause after this call too
            time.sleep(5) 
        except Exception as e:
            print(f"  CRITICAL API ERROR for '{meeting_data['title']}': {e}")
            if "429" in str(e) or "RESOURCE_EXHAUSTED" in str(e):
                print("  ⚠️ Quota Exceeded. Pausing script for 60 seconds before trying NEXT meeting...")
                time.sleep(60)
                # Skip this meeting, try the next one
                continue
            else:
                # If it's another error, try to continue with unknown brand
                brand_details = {"brand_name": "Unknown Brand", "industry": "Unknown"}
        # --- END SAFETY BLOCK ---

        # Step 6: Handle ambiguous result from the LLM
        if brand_details['brand_name'] == 'Unknown Brand' or brand_details['brand_name'].lower() == 'unknown':
            print(f"  Event '{meeting_data['title']}': Title is ambiguous for brand extraction by LLM.")
            # Your notification logic for ambiguous titles can go here if needed.
            # Example:
            # ambiguous_body_html = f"..."
            # send_notification_email(...)
            save_processed_event_id(event_id)
            tag_event_as_processed(calendar_service, event_id) # Tag it so we don't retry
            # Updating the unknown brand name and industry in the master sheet
            index_of_event = updated_meeting_ids.index([event_id]) + 2 # +2 because A1 is header and A2 is first data row
            print(f"  Updating master sheet for event ID '{event_id}' at row {index_of_event} with brand 'Unknown")
            update_values = [[brand_details['brand_name'], brand_details['industry']]]
            body = {
            "valueInputOption": 'USER_ENTERED',  # Use USER_ENTERED to allow date formatting
            "data": [
                {"range": f"Meeting_data!F{index_of_event}:G{index_of_event}", "values": update_values},
                {"range": f"Audit_and_Training!F{index_of_event}:G{index_of_event}", "values": update_values},
                ],
            }
            try:
                resp = (
                    sheets_service.spreadsheets()
                    .values()
                    .batchUpdate(spreadsheetId=master_sheet_id, body=body)
                    .execute()
                    )
                print(f"  Master sheet updated successfully for event ID '{event_id}'.")
                try:
                    print(f" Resetting flag to TRUE for updating owner's sheet for event ID '{event_id}'")
                    values = [["TRUE"]]
                    body = {
                        "valueInputOption": 'USER_ENTERED',  # Use USER_ENTERED to allow date formatting
                        "data": [
                            {"range": f"Meeting_data!{column_index_master['Owner sheet to be updated']}{index_of_event}:{column_index_master['Owner sheet to be updated']}{index_of_event}", "values": values},
                            {"range": f"Audit_and_Training!{column_index_audit['Owner sheet to be updated']}{index_of_event}:{column_index_audit['Owner sheet to be updated']}{index_of_event}", "values": values},
                            ],
                            }
                    resp = (
                        sheets_service.spreadsheets()
                        .values()
                        .batchUpdate(spreadsheetId=master_sheet_id, body=body)
                        .execute()
                        )
                    print(f"  Owner sheet flag reset for '{event_id}'.")
                except HttpError as error:
                    print(f"  Error updating master sheet for event ID '{event_id}': {error}")
            except HttpError as error:
                print(f"  Error updating master sheet for event ID '{event_id}': {error}")
            
            continue

        # Step 7: Merge the successful LLM results into the main meeting_data dictionary
        meeting_data.update(brand_details)

        # ========== NEW CODE STARTS HERE ==========
        # Get LinkedIn profiles for brand attendees
        print(f"  📱 Fetching LinkedIn profiles for brand attendees...")
        brand_attendees_with_linkedin = get_brand_attendees_linkedin_info(
            meeting_data.get('brand_attendees_info', []),
            meeting_data['brand_name'],
            gemini_llm_client
        )
        
        # Replace the old brand attendees info with the new one that has LinkedIn URLs
        meeting_data['brand_attendees_info'] = brand_attendees_with_linkedin
        # ========== NEW CODE ENDS HERE ==========

        # ========== NEW CODE FOR KEY CONTACTS STARTS HERE ==========
        # Find potential key contacts (people NOT in the meeting)
        print(f"  🎯 Finding potential key contacts at {meeting_data['brand_name']}...")
        potential_key_contacts = find_potential_key_contacts(
            meeting_data['brand_name'],
            gemini_llm_client
        )
        
        # Add to meeting data
        meeting_data['potential_key_contacts'] = potential_key_contacts
        # ========== NEW CODE FOR KEY CONTACTS ENDS HERE ==========

        current_brand_name_for_meeting = meeting_data['brand_name']
        target_brand_industry = meeting_data['industry']

        # Updating the brand name and industry in the master sheet
        index_of_event = updated_meeting_ids.index([event_id]) + 2 # +2 because A1 is header and A2 is first data row
        print(f"  Updating master sheet for event ID '{event_id}' at row {index_of_event} with brand '{current_brand_name_for_meeting}' and industry '{target_brand_industry}'")
        update_values = [[current_brand_name_for_meeting, target_brand_industry]]
        body = {
            "valueInputOption": 'USER_ENTERED',  # Use USER_ENTERED to allow date formatting
            "data": [
                {"range": f"Meeting_data!F{index_of_event}:G{index_of_event}", "values": update_values},
                {"range": f"Audit_and_Training!F{index_of_event}:G{index_of_event}", "values": update_values},
                ],
            }
        try:
            resp = (
                sheets_service.spreadsheets()
                .values()
                .batchUpdate(spreadsheetId=master_sheet_id, body=body)
                .execute()
                )
            print(f"  Master sheet updated successfully for event ID '{event_id}'.")
            try:
                print(f" Resetting flag to TRUE for updating owner's sheet for event ID '{event_id}'")
                values = [["TRUE"]]
                body = {
                    "valueInputOption": 'USER_ENTERED',  # Use USER_ENTERED to allow date formatting
                    "data": [
                        {"range": f"Meeting_data!{column_index_master['Owner sheet to be updated']}{index_of_event}:{column_index_master['Owner sheet to be updated']}{index_of_event}", "values": values},
                        {"range": f"Audit_and_Training!{column_index_audit['Owner sheet to be updated']}{index_of_event}:{column_index_audit['Owner sheet to be updated']}{index_of_event}", "values": values},
                        ],
                        }
                resp = (
                    sheets_service.spreadsheets()
                    .values()
                    .batchUpdate(spreadsheetId=master_sheet_id, body=body)
                    .execute()
                    )
                print(f"  Owner sheet flag reset for event ID '{event_id}'.")
            
            except HttpError as error:
                print(f"  Error while resetting flag for '{event_id}': {error}")
        except HttpError as error:
            print(f"  Error updating master sheet for event ID '{event_id}': {error}")

        
        print(f"  LLM identified Brand: '{current_brand_name_for_meeting}', Industry: '{target_brand_industry}'")

        # --- THIS IS THE CORRECTED AND SIMPLIFIED BLOCK ---
        
        # Step 1: Check if necessary services are available.
        if drive_service and sheets_service:
            # If services are available, call the function to get the real data.
            internal_data_result = get_internal_nbh_data_for_brand(
                drive_service=drive_service,
                sheets_service=sheets_service,
                gemini_llm_client=gemini_llm_client,
                current_target_brand_name=current_brand_name_for_meeting,
                target_brand_industry=target_brand_industry, 
                current_meeting_data=meeting_data,
                EXCLUDED_NBH_PSEUDO_NAMES_FOR_FOLLOWUP=EXCLUDED_NBH_PSEUDO_NAMES_FOR_FOLLOWUP,
                AGENT_EMAIL=AGENT_EMAIL,
                master_sheet_id=master_sheet_id 
            )
        else:
            # If services are NOT available, create the default/fallback structure.
            print(f"  Drive/Sheets service not available. Skipping internal data fetch for '{current_brand_name_for_meeting}'.")
            internal_data_result = {
                "llm_summary_string": "Internal NBH Data: Not fetched due to Drive/Sheets service issues.",
                "is_overall_direct_follow_up": False,
                "has_previous_interactions": False,
                "condensed_past_meetings_for_alert": []
            }
        
        # Step 2: Extract the summary string for the LLM brief from the result (either real or default).
        internal_nbh_data_for_brand_str = internal_data_result["llm_summary_string"]
        
        # Step 3: Now, use the result for the leadership alert logic.
        has_prev_interactions_in_main = internal_data_result.get("has_previous_interactions", False)
        is_overall_follow_up_in_main = internal_data_result.get("is_overall_direct_follow_up", False)
        
        
        
                
        # --- >>> LEADERSHIP ALERT LOGIC (FINAL, CORRECTED VERSION) <<< ---

        # Step 1: Extract the flags and data we need from the internal data check.
        is_direct_follow_up = internal_data_result.get("is_overall_direct_follow_up", False)
        has_other_interactions = internal_data_result.get("has_other_past_interactions", False)
        condensed_meetings_for_alert = internal_data_result.get("condensed_past_meetings_for_alert", [])

        # Helper variables for the email body
        upcoming_meeting_title = meeting_data.get('title', 'N/A')
        upcoming_nbh_attendees_list = [att['name'] for att in meeting_data.get('nbh_attendees', [])]
        upcoming_nbh_attendees_str = ", ".join(upcoming_nbh_attendees_list) if upcoming_nbh_attendees_list else "N/A"


        # SCENARIO 1: "Hybrid" Engagement - A follow-up, but other separate threads also exist.
        if is_direct_follow_up and has_other_interactions:
            print("DEBUG: HYBRID SCENARIO DETECTED. Sending nuanced leadership alert.")
            
            alert_subject = f"FYI: Complex Engagement with {current_brand_name_for_meeting} (Follow-up & Separate Threads)"
            
            alert_body_html = f"""
            <html><head><style> body {{ font-family: Arial, sans-serif; }} li {{ margin-bottom: 8px; }} </style></head>
            <body>
                <p>Hello Leadership Team,</p>
                <p>A new meeting has been scheduled with <b>{current_brand_name_for_meeting}</b>. This engagement is complex and requires coordination:</p>
                <ul style="list-style-type:square;">
                    <li>It appears to be a <b>direct follow-up</b> to some recent discussions.</li>
                    <li>However, there are also <b>other, separate historical interactions</b> with this brand.</li>
                </ul>
                <p><b>Upcoming Meeting Details:</b></p>
                <ul>
                    <li><b>Title:</b> {upcoming_meeting_title}</li>
                    <li><b>NBH Attendees:</b> {upcoming_nbh_attendees_str}</li>
                </ul>
                <p>This highlights a need for internal coordination. Context on the separate past interactions is below for awareness:</p>
                <ul>
            """
            if condensed_meetings_for_alert:
                for past_mtg in condensed_meetings_for_alert:
                    alert_body_html += f"<li><b>{past_mtg['date']}:</b> {past_mtg['discussion_summary']} (NBH Team: {past_mtg['nbh_team']})</li>"
            alert_body_html += "</ul><p>Best regards,<br>NBH Meeting Prep Agent</p></body></html>"
            
            # --- CORRECT EMAIL SENDING LOGIC ---
            if gmail_service and leadership_emails:
                email_message = create_email_message(
                    sender=AGENT_EMAIL,
                    to_emails_list=leadership_emails,
                    subject=alert_subject,
                    message_text_html=alert_body_html
                )
                send_gmail_message(gmail_service, 'me', email_message)
                print(f"    Leadership alert for HYBRID scenario with {current_brand_name_for_meeting} sent.")
            else:
                print(f"    WARNING: Leadership alert for {current_brand_name_for_meeting} NOT sent (Gmail service or recipient list unavailable).")


        # SCENARIO 2: "Purely Separate" Engagement - Not a follow-up, but other past interactions exist.
        elif has_other_interactions and not is_direct_follow_up:
            print("DEBUG: PURELY SEPARATE THREAD DETECTED. Sending standard leadership alert.")
            
            alert_subject = f"FYI: New Meeting Scheduled with Existing Brand - {current_brand_name_for_meeting}"
            
            alert_body_html = f"""
            <html><head><style> body {{ font-family: Arial, sans-serif; }} li {{ margin-bottom: 8px; }} </style></head>
            <body>
                <p>Hello Leadership Team,</p>
                <p>A new meeting has been scheduled with <b>{current_brand_name_for_meeting}</b>. This meeting does <b>NOT</b> appear to be a direct follow-up to recent discussions.</p>
                <p>This could indicate a new opportunity or a new NBH team engaging with the client.</p>
                <p><b>Upcoming Meeting Details:</b></p>
                <ul>
                    <li><b>Title:</b> {upcoming_meeting_title}</li>
                    <li><b>NBH Attendees:</b> {upcoming_nbh_attendees_str}</li>
                </ul>
                <p><b>Summary of Past Interactions (for context):</b></p>
                <ul>
            """
            if condensed_meetings_for_alert:
                for past_mtg in condensed_meetings_for_alert:
                    alert_body_html += f"<li><b>{past_mtg['date']}:</b> {past_mtg['discussion_summary']} (NBH Team: {past_mtg['nbh_team']})</li>"
            alert_body_html += "</ul><p>Best regards,<br>NBH Meeting Prep Agent</p></body></html>"
            
            # --- CORRECT EMAIL SENDING LOGIC ---
            if gmail_service and leadership_emails:
                email_message = create_email_message(
                    sender=AGENT_EMAIL,
                    to_emails_list=leadership_emails,
                    subject=alert_subject,
                    message_text_html=alert_body_html
                )
                send_gmail_message(gmail_service, 'me', email_message)
                print(f"    Leadership alert for SEPARATE THREAD with {current_brand_name_for_meeting} sent.")
            else:
                print(f"    WARNING: Leadership alert for {current_brand_name_for_meeting} NOT sent (Gmail service or recipient list unavailable).")

        else:
            # This covers the "clean" cases: a brand-new meeting or a simple follow-up with no other threads.
            print("DEBUG: No leadership alert needed (Clean new meeting or simple follow-up).")
    


        if not gemini_llm_client:
            print(f"  Skipping brief generation for '{meeting_data['title']}': Gemini LLM not available.")
            # Don't mark as processed yet, maybe LLM will be available next run
            continue
        
        if not meeting_data.get('nbh_attendees'): # Check if any NBH humans are there
            print(f"  Event '{meeting_data['title']}': No NBH attendees (other than brandvmeet) to send brief to.")
            save_processed_event_id(event_id)
            tag_event_as_processed(calendar_service, event_id)
            continue


        print(f"  Proceeding with brief generation for: {meeting_data['brand_name']}")
        generated_brief = generate_brief_with_gemini(gemini_llm_client, YOUR_DETAILED_PROMPT_TEMPLATE_GEMINI, meeting_data, internal_nbh_data_for_brand_str)

        FEEDBACK_FORM_URL = "https://forms.gle/Ho9XLKsuGYhWBrBw7"

        # 2. Define the Footer Text (Using Markdown for the email)
        # We use the tagline from your form to keep it consistent.
        feedback_footer = f"""
\n\n
---
Give Your Feedback on The Pre Meeting Briefs. 
👉 Click Here to Fill the Feedback Form ({FEEDBACK_FORM_URL})

"""

        # 3. Append the footer to the generated brief
        # Only add it if the brief was generated successfully (no errors)
        if generated_brief and "Error:" not in generated_brief:
            generated_brief += feedback_footer

        if "Error:" in generated_brief or not generated_brief.strip(): # Check for errors from LLM
            print(f"  Failed to generate brief for '{meeting_data['title']}': {generated_brief}")
            error_body_html = f"""
            <html><body><p>The pre-meeting brief agent encountered an error while generating the brief for:</p>
            <p><b>Event:</b> {meeting_data['title']}<br>
            <b>Brand:</b> {meeting_data['brand_name']}<br>
            <b>Scheduled:</b> {meeting_data['start_time_str']}</p>
            <p><b>Error details:</b> {generated_brief}</p></body></html>"""
            send_notification_email(gmail_service,
                                    f"Error Generating Brief: {meeting_data['title']}",
                                    error_body_html)
            # Don't tag as fully processed if LLM fails, maybe it's temporary.
            # Or use a different tag like [NBH_BRIEF_AGENT_ERROR_V1]
        else:
            print(f"  Successfully generated brief for '{meeting_data['title']}'.")
            send_brief_email(gmail_service, meeting_data, generated_brief)
            tag_event_as_processed(calendar_service, event_id) # Tag only on full success
            set_one_hour_email_reminder(calendar_service, event_id) # Add this call
            save_processed_event_id(event_id)
            # --- Create Google Doc for the brief ---
            BRIEF_FOLDER_ID = "1RhhsFq5NGC2QtHPj8FQaU5BfhxJR5R6I"
            doc_id = create_google_doc_in_folder(
                drive_service,
                BRIEF_FOLDER_ID,
                f"Pre-Meeting Brief - {meeting_data['brand_name']} - {meeting_data['title']}"
            )
            if doc_id:
                write_into_doc(docs_service, doc_id=doc_id, text=generated_brief)
                # Updating doc link in master sheet
                index_of_event = updated_meeting_ids.index([event_id]) + 2 # +2 because A1 is header and A2 is first data row
                update_values = [[f"https://docs.google.com/document/d/{doc_id}"]]
                try:
                    body = {
                        "valueInputOption": 'USER_ENTERED',  # Use USER_ENTERED to allow date formatting
                        "data": [
                            {"range": f"Meeting_data!H{index_of_event}:H{index_of_event}", "values": update_values},
                            {"range": f"Audit_and_Training!H{index_of_event}:H{index_of_event}", "values": update_values},
                            ],
                        }
                    resp = (
                        sheets_service.spreadsheets()
                        .values()
                        .batchUpdate(spreadsheetId=master_sheet_id, body=body)
                        .execute()
                        )
                    print(f"  Master sheet updated with Google Doc link for event ID '{event_id}'.")
                    try:
                        print(f" Resetting flag to TRUE for updating owner's sheet for event ID '{event_id}'")
                        values = [["TRUE"]]
                        body = {
                            "valueInputOption": 'USER_ENTERED',  # Use USER_ENTERED to allow date formatting
                            "data": [
                                {"range": f"Meeting_data!{column_index_master['Owner sheet to be updated']}{index_of_event}:{column_index_master['Owner sheet to be updated']}{index_of_event}", "values": values},
                                {"range": f"Audit_and_Training!{column_index_audit['Owner sheet to be updated']}{index_of_event}:{column_index_audit['Owner sheet to be updated']}{index_of_event}", "values": values},
                                ],
                                }
                        resp = (
                            sheets_service.spreadsheets()
                            .values()
                            .batchUpdate(spreadsheetId=master_sheet_id, body=body)
                            .execute()
                            )
                        print(f"  Owner sheet flag reset for event ID '{event_id}'.")
                    except HttpError as error:
                        print(f"  Error while resetting flag for '{event_id}': {error}")
                except HttpError as error:
                    print(f"  Error updating master sheet with Google Doc link for event ID '{event_id}': {error}")
                # If we have an alternate sheet, update it too
                print(f"  Google Doc created and content written for '{meeting_data['title']}'.")
            
        

    #print(f"Script finished at {datetime.datetime.now()}")

if __name__ == '__main__':
    main()
